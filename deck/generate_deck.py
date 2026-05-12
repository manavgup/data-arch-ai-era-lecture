#!/usr/bin/env python3
"""Generate the 53-slide PPTX deck for 'Data Architecture for the AI Era'.

Usage:
    python deck/generate_deck.py

Output:
    deck/data-architecture-ai-era.pptx

Design system (aligned with Claude Design lecture template):
    Canvas  : 16:9 widescreen (13.333" x 7.5")
    Bg      : #F4F2EC (paper) on ALL slides — light theme, non-negotiable
    BgAlt   : #EAE7DF (paper-alt, used for stat/reality-check slides)
    Accent  : #2D4ADE (signal blue)
    Accent2 : #1A2C8C (deep blue, for secondary accents)
    Ink     : #15171A (primary text)
    Ink2    : #3A3F47 (secondary text)
    Slate   : #5C6470 (tertiary / mono labels)
    Slate2  : #8A93A0 (quaternary)
    Rule    : #C9C4B6 (hairlines / borders)
    RuleSoft: #DCD7C8 (softer hairlines)
    Grid    : #E1DCCB (blueprint grid background)
    Warn    : #B85C00 (ochre — errors, warnings)
    Fonts   : Georgia (display/titles), Calibri (body), Consolas (mono/code)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Design tokens ──────────────────────────────────────────────────────────────

PAPER = RGBColor(0xF4, 0xF2, 0xEC)
PAPER_ALT = RGBColor(0xEA, 0xE7, 0xDF)
ACCENT = RGBColor(0x2D, 0x4A, 0xDE)
ACCENT2 = RGBColor(0x1A, 0x2C, 0x8C)
INK = RGBColor(0x15, 0x17, 0x1A)
INK2 = RGBColor(0x3A, 0x3F, 0x47)
SLATE = RGBColor(0x5C, 0x64, 0x70)
SLATE2 = RGBColor(0x8A, 0x93, 0xA0)
RULE = RGBColor(0xC9, 0xC4, 0xB6)
RULE_SOFT = RGBColor(0xDC, 0xD7, 0xC8)
GRID = RGBColor(0xE1, 0xDC, 0xCB)
WARN = RGBColor(0xB8, 0x5C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Ref-arch lane palette (from Claude Design template) ───────────────────────

LANE_PALETTE = {
    "insight": {"bar": RGBColor(0x1A, 0x2C, 0x8C), "chip": RGBColor(0xE4, 0xE7, 0xF7)},
    "motion": {"bar": RGBColor(0x0E, 0x5E, 0x5C), "chip": RGBColor(0xDE, 0xEA, 0xE9)},
    "storage": {"bar": RGBColor(0x2D, 0x4A, 0xDE), "chip": RGBColor(0xE0, 0xE6, 0xFB)},
    "ingestion": {"bar": RGBColor(0x5B, 0x3F, 0xB8), "chip": RGBColor(0xE8, 0xE2, 0xF5)},
    "sources": {"bar": RGBColor(0x3A, 0x3F, 0x47), "chip": RGBColor(0xE5, 0xE3, 0xDA)},
    "governance": {"bar": RGBColor(0xB8, 0x5C, 0x00), "chip": RGBColor(0xF2, 0xE4, 0xD2)},
    "deploy": {"bar": RGBColor(0x15, 0x17, 0x1A), "chip": RGBColor(0xE5, 0xE3, 0xDA)},
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_DISPLAY = "Georgia"
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"

PT_DISPLAY = 48
PT_TITLE = 28
PT_SUBTITLE = 20
PT_BODY = 14
PT_SMALL = 11
PT_MONO = 10

MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
MARGIN_T = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

ASSETS_DIR = Path(__file__).parent / "assets"

# ── Helper functions ───────────────────────────────────────────────────────────


def _set_slide_bg(slide, color: RGBColor = PAPER):
    """Set slide background to a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide(prs: Presentation):
    """Add a blank slide with paper background."""
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide)
    return slide


def _add_run(
    tf,
    text: str,
    *,
    font_name: str = FONT_BODY,
    size: int = PT_BODY,
    color: RGBColor = INK,
    bold: bool = False,
    italic: bool = False,
):
    """Add a run to the first paragraph of a text frame (helper for single-para boxes)."""
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return run


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_name: str = FONT_BODY,
    size: int = PT_BODY,
    color: RGBColor = INK,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    word_wrap: bool = True,
):
    """Add a simple single-paragraph text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    # Remove vertical auto-fit so text doesn't shrink
    bodyPr = tf.paragraphs[0]._p.getparent().find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set(
            "anchor",
            {
                MSO_ANCHOR.TOP: "t",
                MSO_ANCHOR.MIDDLE: "ctr",
                MSO_ANCHOR.BOTTOM: "b",
            }.get(anchor, "t"),
        )

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_multiline_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list[dict],
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    line_spacing: float | None = None,
):
    """Add a text box with multiple paragraphs.

    Each item in *lines* is a dict:
        text, font_name, size, color, bold, italic, space_after (optional Pt value)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = spec.get("text", "")
        run.font.name = spec.get("font_name", FONT_BODY)
        run.font.size = Pt(spec.get("size", PT_BODY))
        run.font.color.rgb = spec.get("color", INK)
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        if "space_after" in spec:
            p.space_after = Pt(spec["space_after"])
        if "space_before" in spec:
            p.space_before = Pt(spec["space_before"])
        if line_spacing is not None:
            p.line_spacing = Pt(line_spacing)

    return txBox


def add_rule(slide, left, top, width, *, color: RGBColor = RULE, weight: float = 1.5):
    """Add a horizontal rule (line shape)."""
    line = slide.shapes.add_connector(1, left, top, left + width, top)  # MSO_CONNECTOR.STRAIGHT = 1
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_rect(slide, left, top, width, height, *, fill: RGBColor, line_color: RGBColor | None = None):
    """Add a filled rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()  # no border
    return shape


def set_notes(slide, text: str):
    """Set speaker notes on a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_footer(slide, chapter: str, page: int, total: int = 55):
    """Add a standard footer bar: chapter label left, page number right."""
    y = SLIDE_H - Inches(0.45)
    add_rule(slide, MARGIN_L, y, CONTENT_W, color=RULE, weight=0.75)
    y_text = y + Inches(0.08)
    add_textbox(
        slide,
        MARGIN_L,
        y_text,
        Inches(6),
        Inches(0.3),
        chapter,
        font_name=FONT_MONO,
        size=PT_MONO,
        color=SLATE,
    )
    add_textbox(
        slide,
        SLIDE_W - MARGIN_R - Inches(1.5),
        y_text,
        Inches(1.5),
        Inches(0.3),
        f"{page} / {total}",
        font_name=FONT_MONO,
        size=PT_MONO,
        color=SLATE,
        align=PP_ALIGN.RIGHT,
    )


def add_field_guide_header(slide, section_label: str, title: str, *, page: int, chapter: str):
    """Add the standard Field Guide slide header (kicker + title) and footer."""
    # Section kicker
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.5),
        CONTENT_W,
        Inches(0.35),
        section_label,
        font_name=FONT_MONO,
        size=PT_MONO,
        color=ACCENT,
        bold=False,
    )
    # Title
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.9),
        CONTENT_W,
        Inches(0.6),
        title,
        font_name=FONT_DISPLAY,
        size=PT_TITLE,
        color=INK,
        bold=True,
    )
    # Rule below title
    add_rule(slide, MARGIN_L, Inches(1.55), CONTENT_W, color=RULE, weight=1.0)
    # Footer
    add_footer(slide, chapter, page)


def add_table(
    slide,
    left,
    top,
    width,
    data: list[list[str]],
    col_widths: list[float] | None = None,
    *,
    header_font: str = FONT_MONO,
    body_font: str = FONT_BODY,
    header_size: int = PT_SMALL,
    body_size: int = PT_BODY,
    accent_col: int | None = None,
):
    """Add a formatted table.  First row is treated as header."""
    rows = len(data)
    cols = len(data[0]) if data else 0
    row_height = Inches(0.45)
    table_h = row_height * rows

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, table_h)
    tbl = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = ""
            # Remove default paragraph, add formatted one
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = cell_text

            if r_idx == 0:
                # Header row styling
                run.font.name = header_font
                run.font.size = Pt(header_size)
                run.font.color.rgb = SLATE
                run.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = PAPER
            else:
                # Body row styling
                run.font.name = body_font
                run.font.size = Pt(body_size)
                if accent_col is not None and c_idx == accent_col:
                    run.font.color.rgb = ACCENT
                    run.font.bold = True
                else:
                    run.font.color.rgb = INK

                cell.fill.solid()
                cell.fill.fore_color.rgb = PAPER

            # Cell borders — thin rule colour
            _set_cell_border(cell, RULE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p.alignment = PP_ALIGN.LEFT

    return table_shape


def _set_cell_border(cell, color: RGBColor, width: str = "6350"):
    """Set thin borders on a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = tcPr.find(qn(edge))
        if ln is None:
            ln = tcPr.makeelement(qn(edge), {})
            tcPr.append(ln)
        ln.set("w", width)
        solidFill = ln.find(qn("a:solidFill"))
        if solidFill is None:
            solidFill = ln.makeelement(qn("a:solidFill"), {})
            ln.append(solidFill)
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = solidFill.makeelement(qn("a:srgbClr"), {})
            solidFill.append(srgb)
        srgb.set("val", str(color))


def add_big_quote(slide, quote_text: str, attribution: str, *, page: int, chapter: str):
    """Add a big-quote layout slide."""
    # Large quote mark
    add_textbox(
        slide,
        Inches(1.2),
        Inches(1.5),
        Inches(1.5),
        Inches(1.5),
        "\u201c",
        font_name=FONT_DISPLAY,
        size=96,
        color=ACCENT,
        bold=True,
    )
    # Quote text
    add_textbox(
        slide,
        Inches(2.4),
        Inches(2.0),
        Inches(9.0),
        Inches(3.0),
        quote_text,
        font_name=FONT_DISPLAY,
        size=PT_TITLE,
        color=INK,
        italic=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Attribution
    if attribution:
        add_textbox(
            slide,
            Inches(2.4),
            Inches(5.2),
            Inches(9.0),
            Inches(0.4),
            attribution,
            font_name=FONT_MONO,
            size=PT_MONO,
            color=SLATE,
        )
    add_footer(slide, chapter, page)


def add_section_divider(slide, block_label: str, title: str, blurb: str, duration: str, *, page: int):
    """Add a section divider slide."""
    # Accent blue header bar at top
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.4), fill=ACCENT)

    # Block label
    add_textbox(
        slide,
        MARGIN_L,
        Inches(2.0),
        CONTENT_W,
        Inches(0.4),
        block_label,
        font_name=FONT_MONO,
        size=PT_MONO,
        color=SLATE,
    )
    # Title
    add_textbox(
        slide,
        MARGIN_L,
        Inches(2.6),
        CONTENT_W,
        Inches(1.0),
        title,
        font_name=FONT_DISPLAY,
        size=36,
        color=INK,
        bold=True,
    )
    # Blurb
    add_textbox(
        slide,
        MARGIN_L,
        Inches(3.8),
        Inches(10),
        Inches(1.2),
        blurb,
        font_name=FONT_DISPLAY,
        size=18,
        color=INK2,
        italic=True,
    )
    # Bottom timing info
    add_textbox(
        slide,
        MARGIN_L,
        SLIDE_H - Inches(0.8),
        Inches(4),
        Inches(0.4),
        duration,
        font_name=FONT_MONO,
        size=PT_MONO,
        color=SLATE,
    )
    add_footer(slide, title.rstrip("."), page)


def add_diagram_or_placeholder(slide, filename: str, label: str, top=None, height=None):
    """Try to place a PNG from assets as full-slide image. If missing, add placeholder.

    Args:
        top: vertical offset for the image (default Inches(0) = full bleed).
        height: image height (default SLIDE_H). When top is set, defaults to
                SLIDE_H - top so the diagram fills below the header zone.
    """
    img_top = top if top is not None else Inches(0)
    img_h = height if height is not None else (SLIDE_H - img_top if top is not None else SLIDE_H)
    img_path = ASSETS_DIR / filename
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path),
            Inches(0),
            img_top,
            SLIDE_W,
            img_h,
        )
    else:
        placeholder_top = img_top + Inches(0.3) if top is not None else Inches(1.5)
        # Placeholder rectangle
        add_rect(
            slide,
            Inches(1.5),
            placeholder_top,
            SLIDE_W - Inches(3),
            img_h - Inches(1),
            fill=RGBColor(0xE8, 0xE5, 0xDB),
            line_color=RULE,
        )
        add_textbox(
            slide,
            Inches(2),
            placeholder_top + Inches(1),
            SLIDE_W - Inches(4),
            Inches(1.5),
            f"DIAGRAM: {label}",
            font_name=FONT_MONO,
            size=PT_SUBTITLE,
            color=SLATE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def add_bullet_list(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    *,
    font_name: str = FONT_BODY,
    size: int = PT_BODY,
    color: RGBColor = INK,
    bullet_char: str = "\u2022",
    line_spacing: float | None = None,
):
    """Add a bulleted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        p.space_after = Pt(4)

    return txBox


# ── Slide builders ─────────────────────────────────────────────────────────────
# Each function creates exactly one slide (or a small group for 08b-08e).


def make_slide_01_cover(prs):
    """Slide 01 -- Title / Cover."""
    slide = add_slide(prs)

    # Top bar
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.4),
        CONTENT_W,
        Inches(0.3),
        "A FIELD GUIDE \u00b7 V1.0 \u00b7 2026    |    HALF-DAY \u00b7 LECTURE + HANDS-ON",
        font_name=FONT_MONO,
        size=PT_MONO,
        color=SLATE,
    )

    # Kicker
    add_textbox(
        slide,
        MARGIN_L,
        Inches(1.4),
        CONTENT_W,
        Inches(0.4),
        "VOLUME ONE \u2014 FOUNDATIONS",
        font_name=FONT_MONO,
        size=PT_SMALL,
        color=ACCENT,
        bold=True,
    )

    # Title (two lines)
    add_textbox(
        slide,
        MARGIN_L,
        Inches(2.0),
        CONTENT_W,
        Inches(2.0),
        "Data Architecture\nfor the AI Era.",
        font_name=FONT_DISPLAY,
        size=PT_DISPLAY,
        color=INK,
        bold=True,
    )

    # Rule
    add_rule(slide, MARGIN_L, Inches(4.2), Inches(10), color=RULE, weight=2.0)

    # Subtitle (italic)
    add_textbox(
        slide,
        MARGIN_L,
        Inches(4.5),
        Inches(10),
        Inches(1.0),
        "A field guide for sellers, architects, and SSRs \u2014 six patterns, one BFSI scenario, six runnable notebooks.",
        font_name=FONT_DISPLAY,
        size=PT_SUBTITLE,
        color=INK2,
        italic=True,
    )

    # Bottom rule
    add_rule(slide, MARGIN_L, Inches(6.3), CONTENT_W, color=RULE, weight=1.0)

    # Bottom bar
    add_textbox(
        slide,
        MARGIN_L,
        Inches(6.5),
        CONTENT_W,
        Inches(0.3),
        "MAPLE TRUST BANK \u00b7 CASE STUDY    |    MANAV GUPTA    |    VP & CTO, TECHNICAL SALES",
        font_name=FONT_MONO,
        size=PT_MONO,
        color=SLATE,
    )

    set_notes(
        slide,
        (
            "Welcome everyone. This is a working session, not a PowerPoint marathon. "
            "You have six Jupyter notebooks in front of you. By the end of today, you'll know "
            "which data architecture pattern to propose for any customer conversation \u2014 and "
            "you'll have run the code to prove it."
        ),
    )


def make_slide_02_hook(prs):
    """Slide 02 -- The Hook (Big Quote)."""
    slide = add_slide(prs)
    add_big_quote(
        slide,
        "The data architecture you designed in 2019 is what\u2019s breaking your AI strategy in 2026.",
        "",
        page=2,
        chapter="Opening",
    )
    set_notes(
        slide,
        (
            'This is the thesis. Not "your architecture is bad" \u2014 it was right for 2019. '
            "But AI changed the read patterns, the latency budgets, and what counts as data. "
            "Today we fix the gap."
        ),
    )


def make_slide_03_refarch(prs):
    """Slide 03 -- Full Reference Architecture diagram."""
    slide = add_slide(prs)

    # Title above diagram
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.3),
        CONTENT_W,
        Inches(0.5),
        "Reference Architecture \u2014 IBM Software Hub (Cloud Pak for Data Platform)",
        font_name=FONT_DISPLAY,
        size=PT_SUBTITLE,
        color=INK,
        bold=True,
    )
    add_diagram_or_placeholder(
        slide, "refarch-full.png", "Full Reference Architecture \u2014 IBM Software Hub", top=Inches(1.0)
    )

    add_footer(slide, "Opening", 3)

    set_notes(
        slide,
        (
            "This is the slide everyone has seen 400 times. Today is the day someone explains it to you. "
            "Every block of this lecture will light up different swimlanes. By the end, you'll have the "
            "whole diagram populated in your head \u2014 and you'll know which swimlane to point at for "
            "any customer conversation. Let me walk you through it at 30,000 feet."
        ),
    )


def _add_lane_block(slide, left, top, width, height, label, items, bar_color, chip_bg, chip_fg):
    """Draw a single swimlane: colored header bar + chip labels inside."""
    # Header bar
    add_rect(slide, left, top, width, Inches(0.3), fill=bar_color)
    add_textbox(
        slide,
        left + Inches(0.12),
        top + Inches(0.02),
        width - Inches(0.24),
        Inches(0.26),
        label,
        font_name=FONT_MONO,
        size=8,
        color=WHITE,
        bold=False,
    )
    # Chip area
    chip_y = top + Inches(0.34)
    chip_x = left + Inches(0.12)
    row_h = Inches(0.28)
    col_gap = Inches(0.08)
    max_x = left + width - Inches(0.12)

    for item in items:
        text_w = Inches(max(1.2, len(item) * 0.085))
        if chip_x + text_w > max_x:
            chip_x = left + Inches(0.12)
            chip_y += row_h
        add_rect(slide, chip_x, chip_y, text_w, Inches(0.24), fill=chip_bg, line_color=RULE)
        add_textbox(
            slide,
            chip_x + Inches(0.06),
            chip_y + Inches(0.02),
            text_w - Inches(0.12),
            Inches(0.2),
            item,
            font_name=FONT_BODY,
            size=8,
            color=chip_fg,
        )
        chip_x += text_w + col_gap


def make_slide_03b_refarch_overview(prs):
    """Slide 03b -- Ref Arch Overview: color-coded swimlane diagram."""
    slide = add_slide(prs)

    # Header
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.25),
        CONTENT_W,
        Inches(0.25),
        "\u00a7 00 \u2014 SOFTWARE HUB \u00b7 REFERENCE ARCHITECTURE",
        font_name=FONT_MONO,
        size=8,
        color=ACCENT,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.52),
        CONTENT_W,
        Inches(0.5),
        "The diagram every Software Hub conversation comes back to.",
        font_name=FONT_DISPLAY,
        size=24,
        color=INK,
        bold=True,
    )

    # Layout geometry
    gov_w = Inches(2.0)  # left governance rail
    gov_left = MARGIN_L
    lanes_left = gov_left + gov_w + Inches(0.1)
    lanes_w = CONTENT_W - gov_w - Inches(0.1)
    base_y = Inches(1.2)
    lane_h = Inches(0.72)
    lane_gap = Inches(0.06)

    # ── Left rail: Governance ──
    gov_bar = LANE_PALETTE["governance"]["bar"]
    gov_chip_bg = LANE_PALETTE["governance"]["chip"]
    add_rect(slide, gov_left, base_y, gov_w, Inches(0.3), fill=gov_bar)
    add_textbox(
        slide,
        gov_left + Inches(0.1),
        base_y + Inches(0.02),
        gov_w - Inches(0.2),
        Inches(0.26),
        "INFORMATION & MODEL\nMANAGEMENT & GOVERNANCE",
        font_name=FONT_MONO,
        size=7,
        color=WHITE,
    )
    gov_items = [
        "Business Glossary",
        "Data Lineage",
        "Data Quality",
        "Regulatory Privacy",
        "MDM / Match 360",
        "Catalog",
    ]
    chip_y = base_y + Inches(0.36)
    for item in gov_items:
        add_rect(
            slide,
            gov_left + Inches(0.08),
            chip_y,
            gov_w - Inches(0.16),
            Inches(0.22),
            fill=gov_chip_bg,
            line_color=RULE,
        )
        add_textbox(
            slide,
            gov_left + Inches(0.16),
            chip_y + Inches(0.02),
            gov_w - Inches(0.32),
            Inches(0.18),
            item,
            font_name=FONT_BODY,
            size=8,
            color=LANE_PALETTE["governance"]["bar"],
        )
        chip_y += Inches(0.24)
    # Security + Platform text
    add_textbox(
        slide,
        gov_left + Inches(0.1),
        chip_y + Inches(0.08),
        gov_w - Inches(0.2),
        Inches(0.3),
        "SECURITY\nPLATFORM INFRASTRUCTURE",
        font_name=FONT_MONO,
        size=7,
        color=SLATE,
    )

    # ── Center: 5 stacked lanes ──
    lanes = [
        (
            "ACTIONABLE INSIGHT",
            "insight",
            [
                "Enhanced Applications",
                "Customer Insights",
                "New Business Models",
                "Forecasting, Planning & Analysis",
                "Compliance & Fraud",
                "Security Operations & Risk",
                "Discovery & Exploration",
            ],
        ),
        (
            "ANALYTICS IN-MOTION",
            "motion",
            ["Real-time scoring", "Streaming aggregates", "Event triggers", "Online features"],
        ),
        (
            "ANALYTICAL DATA MANAGEMENT & STORAGE",
            "storage",
            ["Warehouse", "Lakehouse", "Lake (Object Storage)", "Vector / Search", "Operational Stores"],
        ),
        (
            "INGESTION & INTEGRATION \u00b7 DATA ACCESS",
            "ingestion",
            ["Batch ETL", "CDC", "Streaming", "API / Files", "Federation", "Replication"],
        ),
        (
            "DATA SOURCES",
            "sources",
            [
                "System of Record",
                "Application Data",
                "Transactional",
                "Third-Party",
                "Social",
                "Weather",
                "Internet Data Sets",
                "Content Services",
                "Image & Video",
                "Machine & Sensor",
            ],
        ),
    ]

    y = base_y
    for label, lane_key, items in lanes:
        lp = LANE_PALETTE[lane_key]
        _add_lane_block(slide, lanes_left, y, lanes_w, lane_h, label, items, lp["bar"], lp["chip"], lp["bar"])
        y += lane_h + lane_gap

    # Footer annotations
    add_textbox(
        slide,
        MARGIN_L,
        y + Inches(0.1),
        Inches(7),
        Inches(0.2),
        "ADAPTED FROM SOFTWARE HUB 5.2 REFERENCE ARCHITECTURE \u00b7 FOR TEACHING PURPOSES",
        font_name=FONT_MONO,
        size=7,
        color=SLATE,
    )
    add_textbox(
        slide,
        SLIDE_W - MARGIN_R - Inches(5),
        y + Inches(0.1),
        Inches(5),
        Inches(0.2),
        "EACH PATTERN IN THIS DECK LIVES IN 1\u20132 LANES ABOVE",
        font_name=FONT_MONO,
        size=7,
        color=SLATE,
        align=PP_ALIGN.RIGHT,
    )
    add_footer(slide, "Reference Architecture", 3)

    set_notes(
        slide,
        (
            "This is the reference architecture \u2014 the one diagram the whole half-day comes back to. "
            "Five horizontal lanes: sources at the bottom, ingestion, then analytical storage in the center "
            "(where watsonx.data lives), analytics in-motion above that, and actionable insight at the top. "
            "Governance runs vertically on the left. Every pattern we discuss maps to one or two of these lanes."
        ),
    )


def make_slide_03c_refarch_products(prs):
    """Slide 03c -- Ref Arch Products: IBM components named per lane."""
    slide = add_slide(prs)

    # Header
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.25),
        CONTENT_W,
        Inches(0.25),
        "\u00a7 00 \u2014 SOFTWARE HUB \u00b7 PRODUCT MAPPING",
        font_name=FONT_MONO,
        size=8,
        color=ACCENT,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.52),
        CONTENT_W,
        Inches(0.5),
        "The same five lanes \u2014 now with the IBM components named.",
        font_name=FONT_DISPLAY,
        size=24,
        color=INK,
        bold=True,
    )

    # Layout geometry
    gov_w = Inches(1.7)
    deploy_w = Inches(1.7)
    gov_left = MARGIN_L
    deploy_left = SLIDE_W - MARGIN_R - deploy_w
    lanes_left = gov_left + gov_w + Inches(0.08)
    lanes_w = deploy_left - lanes_left - Inches(0.08)
    base_y = Inches(1.15)
    lane_h = Inches(0.62)
    lane_gap = Inches(0.05)

    # ── Left rail: IBM Knowledge Catalog ──
    gov_bar = LANE_PALETTE["governance"]["bar"]
    gov_chip_bg = LANE_PALETTE["governance"]["chip"]
    add_rect(slide, gov_left, base_y, gov_w, Inches(0.25), fill=gov_bar)
    add_textbox(
        slide,
        gov_left + Inches(0.08),
        base_y + Inches(0.02),
        gov_w - Inches(0.16),
        Inches(0.21),
        "IBM KNOWLEDGE CATALOG",
        font_name=FONT_MONO,
        size=7,
        color=WHITE,
    )
    gov_items = [
        "Business Glossary",
        "Data Lineage",
        "Metadata Enrich",
        "Regulatory Governance",
        "Data Quality",
        "Product Inventory",
        "AI Accelerators",
        "MDM Match 360",
        "Privacy",
        "Data Master",
        "watsonx.ai Factsheets",
    ]
    chip_y = base_y + Inches(0.3)
    for item in gov_items:
        add_rect(
            slide,
            gov_left + Inches(0.06),
            chip_y,
            gov_w - Inches(0.12),
            Inches(0.19),
            fill=gov_chip_bg,
            line_color=RULE,
        )
        add_textbox(
            slide,
            gov_left + Inches(0.12),
            chip_y + Inches(0.01),
            gov_w - Inches(0.24),
            Inches(0.17),
            item,
            font_name=FONT_MONO,
            size=7,
            color=gov_bar,
        )
        chip_y += Inches(0.21)
    # Security label
    add_textbox(
        slide,
        gov_left + Inches(0.06),
        chip_y + Inches(0.04),
        gov_w - Inches(0.12),
        Inches(0.3),
        "GUARDIUM DATA PROTECTION\nSECURITY (PRE-INTEGRATED)",
        font_name=FONT_MONO,
        size=6,
        color=SLATE,
    )

    # ── Right rail: Deploy Anywhere ──
    deploy_bar = LANE_PALETTE["deploy"]["bar"]
    add_rect(slide, deploy_left, base_y, deploy_w, Inches(0.25), fill=deploy_bar)
    add_textbox(
        slide,
        deploy_left + Inches(0.08),
        base_y + Inches(0.02),
        deploy_w - Inches(0.16),
        Inches(0.21),
        "DEPLOY ANYWHERE",
        font_name=FONT_MONO,
        size=7,
        color=WHITE,
    )
    deploy_opts = [
        ("IBM Cloud", "SaaS / managed"),
        ("On-Premise", "Customer DC"),
        ("Hyper-converged", "Single-node teaching stack"),
    ]
    dy = base_y + Inches(0.35)
    for title, sub in deploy_opts:
        add_textbox(
            slide,
            deploy_left + Inches(0.1),
            dy,
            deploy_w - Inches(0.2),
            Inches(0.2),
            title,
            font_name=FONT_DISPLAY,
            size=11,
            color=INK,
            bold=True,
        )
        add_textbox(
            slide,
            deploy_left + Inches(0.1),
            dy + Inches(0.18),
            deploy_w - Inches(0.2),
            Inches(0.15),
            sub,
            font_name=FONT_BODY,
            size=7,
            color=SLATE,
        )
        dy += Inches(0.42)
    add_textbox(
        slide,
        deploy_left + Inches(0.1),
        dy + Inches(0.1),
        deploy_w - Inches(0.2),
        Inches(0.2),
        "Same RA. Three substrates.",
        font_name=FONT_MONO,
        size=7,
        color=SLATE,
    )

    # ── Center: 5 stacked lanes with IBM products ──
    lanes = [
        (
            "ACTIONABLE INSIGHT",
            "insight",
            [
                "Watson Studio",
                "watsonx.ai",
                "Watson OpenScale",
                "Watson Machine Learning",
                "Orchestration Pipelines",
                "SPSS Modeler",
                "Decision Optimization",
                "Cognos Dashboards",
                "Cognos Analytics",
                "Data Intelligence",
            ],
        ),
        ("ANALYTICS IN-MOTION", "motion", ["Apache Kafka", "Data Replication", "Presto (connector)", "Connectivity"]),
        (
            "ANALYTICAL STORAGE \u00b7 ON SOFTWARE HUB",
            "storage",
            [
                "watsonx.data",
                "Db2",
                "Db2 Warehouse (SMP, MPP)",
                "MongoDB",
                "EDB PostgreSQL",
                "Informix",
                "Apache Iceberg",
                "Delta Lake",
                "Milvus",
                "Apache Spark SQL",
                "Hadoop Execution",
                "Data Virtualization",
            ],
        ),
        (
            "INGESTION & INTEGRATION \u00b7 DATA ACCESS",
            "ingestion",
            [
                "Data Integration",
                "Apache Spark (Streaming)",
                "Search",
                "Watson Studio",
                "Data Catalog",
                "Data Refinery",
                "Planning Analytics",
            ],
        ),
        (
            "DATA SOURCES",
            "sources",
            [
                "System of Record",
                "Application Data",
                "Content Services",
                "Images & Video",
                "Machine & Sensor",
                "Social",
                "Watson AI Services (WA, WD, Watson Speech)",
            ],
        ),
    ]

    y = base_y
    for label, lane_key, items in lanes:
        lp = LANE_PALETTE[lane_key]
        _add_lane_block(slide, lanes_left, y, lanes_w, lane_h, label, items, lp["bar"], lp["chip"], lp["bar"])
        y += lane_h + lane_gap

    # Footer
    add_textbox(
        slide,
        SLIDE_W - MARGIN_R - Inches(6),
        y + Inches(0.06),
        Inches(6),
        Inches(0.2),
        "ADAPTED FROM SOFTWARE HUB 5.2 \u00b7 CLOUD PAK FOR DATA PLATFORM",
        font_name=FONT_MONO,
        size=7,
        color=SLATE,
        align=PP_ALIGN.RIGHT,
    )
    add_footer(slide, "Reference Architecture", 4)

    set_notes(
        slide,
        (
            "Same five lanes, now with the actual IBM products named. watsonx.data is in the center \u2014 "
            "the analytical storage lane. Presto, Spark, Db2 are query engines. Knowledge Catalog is the "
            "governance rail on the left. The right rail shows this deploys on IBM Cloud, on-premise, or a "
            "hyper-converged single-node stack for teaching. Every product maps to one lane."
        ),
    )


def make_slide_04_roles(prs):
    """Slide 04 -- Why Each Role Cares (3-column)."""
    slide = add_slide(prs)
    add_field_guide_header(slide, "OPENING", "Why each role cares.", page=4, chapter="Opening")

    col_w = Inches(3.6)
    y_top = Inches(1.9)
    gap = Inches(0.3)
    col_h = Inches(3.5)

    roles = [
        ("SELLERS", "Stop selling \u201cAI.\u201d Start selling the data plane that makes AI work."),
        (
            "ARCHITECTS",
            "The integration contract between your enterprise and your agents \u2014 get it wrong and nothing else matters.",
        ),
        ("SSRs", "Every incident is now a governance incident. You need to know which swimlane broke."),
    ]

    for i, (role, desc) in enumerate(roles):
        x = MARGIN_L + (col_w + gap) * i
        # Role label
        add_textbox(
            slide,
            x,
            y_top,
            col_w,
            Inches(0.4),
            role,
            font_name=FONT_MONO,
            size=PT_SMALL,
            color=ACCENT,
            bold=True,
        )
        # Vertical rule below header
        add_rule(slide, x, y_top + Inches(0.45), col_w, color=RULE, weight=0.75)
        # Description
        add_textbox(
            slide,
            x,
            y_top + Inches(0.65),
            col_w,
            col_h,
            desc,
            font_name=FONT_BODY,
            size=PT_BODY,
            color=INK,
        )

    set_notes(
        slide,
        (
            "Ground rules: interrupt freely. We have two live notebooks in Block 4, two notebook cameos "
            "in Block 1, and one architecture critique exercise. If I'm boring you, say so \u2014 this is "
            "a working session."
        ),
    )


# ── Block 1: Foundations ───────────────────────────────────────────────────────


def make_slide_05_divider_foundations(prs):
    """Slide 05 -- Section Divider: Foundations."""
    slide = add_slide(prs)
    add_section_divider(
        slide,
        "BLOCK 1 OF 5",
        "Foundations.",
        "Storage, compute, catalog, governance, observability \u2014 the five primitives. Then the six patterns.",
        "45 MIN",
        page=5,
    )
    set_notes(
        slide,
        (
            "Block 1 is the foundation. If you already know what a lakehouse is, stay with me anyway \u2014 "
            "I'm going to reframe it. If you don't know what a lakehouse is, you will in 45 minutes."
        ),
    )


def make_slide_06_reinvention(prs):
    """Slide 06 -- Why Data Architecture Keeps Getting Reinvented."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 1.1 \u2014 WHY WE KEEP REINVENTING DATA ARCHITECTURE",
        "The 30-year arc.",
        page=6,
        chapter="Block 1 \u2014 Foundations",
    )

    # Timeline (horizontal labels)
    stages = ["OLTP", "Warehouse", "Lake", "Lakehouse", "Mesh", "Fabric"]
    y = Inches(2.8)
    box_w = Inches(1.7)
    gap = Inches(0.25)
    total_w = len(stages) * box_w + (len(stages) - 1) * gap
    start_x = (SLIDE_W - total_w) / 2

    for i, stage in enumerate(stages):
        x = start_x + i * (box_w + gap)
        add_rect(
            slide, x, y, box_w, Inches(0.65), fill=ACCENT if i >= 3 else RGBColor(0xDD, 0xDA, 0xD0), line_color=RULE
        )
        text_color = WHITE if i >= 3 else INK
        add_textbox(
            slide,
            x,
            y + Inches(0.1),
            box_w,
            Inches(0.5),
            stage,
            font_name=FONT_MONO,
            size=PT_SMALL,
            color=text_color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        # Arrow between boxes
        if i < len(stages) - 1:
            add_textbox(
                slide,
                x + box_w,
                y + Inches(0.1),
                gap,  # already in EMU from Inches(0.25)
                Inches(0.5),
                "\u2192",
                font_name=FONT_BODY,
                size=PT_BODY,
                color=SLATE,
                align=PP_ALIGN.CENTER,
            )

    # Narrative below timeline
    add_textbox(
        slide,
        MARGIN_L,
        Inches(4.0),
        CONTENT_W,
        Inches(1.5),
        "Each generation solved the previous generation\u2019s scaling problem and created a new governance problem.\n"
        "The pendulum: centralize \u2192 decentralize \u2192 federate.",
        font_name=FONT_BODY,
        size=PT_BODY,
        color=INK2,
    )

    set_notes(
        slide,
        (
            "Every generation solved the previous generation's scaling problem and created a new governance "
            "problem. Warehouses solved OLTP reporting but couldn't handle unstructured data. Lakes solved "
            "storage cost but created swamps. Lakehouses added ACID but didn't solve org design. Mesh solved "
            "org design but assumed product thinking maturity nobody had. Fabric tries to automate the "
            "governance layer. The cycle continues."
        ),
    )


def make_slide_07_primitives(prs):
    """Slide 07 -- Five Primitives."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 1.2 \u2014 THE FIVE PRIMITIVES",
        "Not four. Five.",
        page=7,
        chapter="Block 1 \u2014 Foundations",
    )

    data = [
        ["PRIMITIVE", "WHAT IT COVERS"],
        ["Storage", "Object stores, columnar formats (Parquet, ORC), open table formats (Iceberg, Delta, Hudi)"],
        ["Compute", "Query engines (Presto/Trino, Spark, DuckDB), separation of storage and compute"],
        ["Catalog", "Hive, Unity, Polaris, Nessie \u2014 catalog wars are the new format wars"],
        ["Governance", "Lineage, access, classification, lifecycle"],
        ["Observability", "Quality, freshness, schema drift, pipeline SLOs, lineage-driven impact analysis"],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(1.8),
        CONTENT_W,
        data,
        col_widths=[2.5, 9.3],
        accent_col=0,
    )

    set_notes(
        slide,
        (
            "Most people teach four primitives and leave out observability. That was fine when data pipelines "
            "ran overnight and you checked them in the morning. It's not fine when an agent loop needs fresh "
            "data in milliseconds. Observability is now first-class \u2014 quality monitoring, freshness SLOs, "
            "schema drift detection. Three layers: data observability, model observability, agent trace "
            "observability. Most orgs have one of the three. Almost nobody has all three."
        ),
    )


def make_slide_08_decoder_ring(prs):
    """Slide 08 -- Pattern Decoder Ring (table)."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 1.3 \u2014 PATTERN DECODER RING",
        "Six patterns. When to use. When each is actively wrong.",
        page=8,
        chapter="Block 1 \u2014 Foundations",
    )

    data = [
        ["PATTERN", "BEST FOR", "ANTI-PATTERN"],
        ["Data Warehouse", "Stable schemas, BI, regulatory reporting", "Unstructured/AI workloads"],
        ["Data Lake", "Cheap raw storage, ML training sets", "Anything needing ACID or governance"],
        ["Lakehouse", "Unified analytics + ML on open formats", "Sub-second OLTP"],
        ["Data Virtualization", "Federated queries, data residency", "High-volume scans, latency-sensitive"],
        ["Data Mesh", "Large orgs with strong domain teams", "Orgs without product thinking"],
        ["Data Fabric", "Heterogeneous estates, federation", "Greenfield single-cloud"],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(1.8),
        CONTENT_W,
        data,
        col_widths=[3.0, 4.9, 3.9],
        accent_col=0,
    )

    set_notes(
        slide,
        (
            "Memorize this table. When a customer describes their problem, you should be able to point at "
            'the right row in under 10 seconds. The anti-pattern column is more important than the "best for" '
            "column \u2014 knowing when NOT to propose a pattern is what separates a seller from a consultant."
        ),
    )


def _make_pattern_swimlane_slide(
    prs,
    title: str,
    subtitle: str,
    highlight_lanes: set[str],
    active_items: dict[str, list[str]] | None,
    page: int,
    notes: str,
):
    """Build a pattern architecture slide as a color-coded swimlane diagram.

    highlight_lanes: set of lane keys to render at full opacity.
    active_items: optional dict mapping lane key -> list of items to render as chips.
                  If None, uses the standard ref-arch items for highlighted lanes.
    """
    slide = add_slide(prs)

    # Header
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.25),
        CONTENT_W,
        Inches(0.25),
        "\u00a7 1.3 \u2014 PATTERN DECODER RING",
        font_name=FONT_MONO,
        size=8,
        color=ACCENT,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.50),
        CONTENT_W,
        Inches(0.35),
        title,
        font_name=FONT_DISPLAY,
        size=22,
        color=INK,
        bold=True,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.88),
        CONTENT_W,
        Inches(0.25),
        subtitle,
        font_name=FONT_BODY,
        size=10,
        color=SLATE,
        italic=True,
    )

    # Standard lane definitions
    STD_LANES = [
        (
            "ACTIONABLE INSIGHT",
            "insight",
            ["Enhanced Applications", "Customer Insights", "Compliance & Fraud", "Discovery & Exploration"],
        ),
        ("ANALYTICS IN-MOTION", "motion", ["Real-time scoring", "Streaming aggregates", "Event triggers"]),
        (
            "ANALYTICAL DATA MANAGEMENT & STORAGE",
            "storage",
            ["Warehouse", "Lakehouse", "Lake (Object Storage)", "Vector / Search"],
        ),
        ("INGESTION & INTEGRATION", "ingestion", ["Batch ETL", "CDC", "Streaming", "API / Files", "Federation"]),
        ("DATA SOURCES", "sources", ["System of Record", "Application Data", "Third-Party", "Content Services"]),
    ]

    # Layout
    lanes_left = MARGIN_L
    lanes_w = CONTENT_W
    base_y = Inches(1.2)
    lane_h = Inches(0.72)
    lane_gap = Inches(0.06)
    dim_color = RGBColor(0xD0, 0xCE, 0xC6)

    y = base_y
    for label, lane_key, default_items in STD_LANES:
        lp = LANE_PALETTE[lane_key]
        is_active = lane_key in highlight_lanes
        items = (active_items or {}).get(lane_key, default_items) if is_active else []

        bar_color = lp["bar"] if is_active else dim_color
        chip_bg = lp["chip"] if is_active else PAPER
        chip_fg = lp["bar"] if is_active else SLATE2

        # Always draw header bar
        add_rect(slide, lanes_left, y, lanes_w, Inches(0.3), fill=bar_color)
        label_color = WHITE if is_active else SLATE
        add_textbox(
            slide,
            lanes_left + Inches(0.12),
            y + Inches(0.02),
            lanes_w - Inches(0.24),
            Inches(0.26),
            label,
            font_name=FONT_MONO,
            size=8,
            color=label_color,
        )

        # Draw chips only for active lanes
        if items:
            chip_y = y + Inches(0.34)
            chip_x = lanes_left + Inches(0.12)
            max_x = lanes_left + lanes_w - Inches(0.12)
            for item in items:
                text_w = Inches(max(1.2, len(item) * 0.085))
                if chip_x + text_w > max_x:
                    chip_x = lanes_left + Inches(0.12)
                    chip_y += Inches(0.28)
                add_rect(slide, chip_x, chip_y, text_w, Inches(0.24), fill=chip_bg, line_color=RULE)
                add_textbox(
                    slide,
                    chip_x + Inches(0.06),
                    chip_y + Inches(0.02),
                    text_w - Inches(0.12),
                    Inches(0.2),
                    item,
                    font_name=FONT_BODY,
                    size=8,
                    color=chip_fg,
                )
                chip_x += text_w + Inches(0.08)

        y += lane_h + lane_gap

    add_footer(slide, "Block 1 \u2014 Foundations", page)
    set_notes(slide, notes)
    return slide


def make_slide_08b_lake(prs):
    """Slide 08b -- Data Lake Architecture (swimlane diagram)."""
    _make_pattern_swimlane_slide(
        prs,
        title="Data Lake",
        subtitle="Cheap object storage holding everything raw. Schema-on-read, when someone reads it.",
        highlight_lanes={"sources", "ingestion", "storage"},
        active_items={
            "sources": ["Core Banking", "CRM", "Policy PDFs", "Card Stream"],
            "ingestion": ["Batch ETL", "CDC", "API / Files"],
            "storage": ["Lake (Object Storage)", "Parquet", "JSON", "CSV", "Raw PDFs"],
        },
        page=9,
        notes=(
            "The lake. Cheap storage for everything. Parquet, JSON, CSV, PDFs \u2014 dump it all in and figure "
            "out the schema later. Some teams did figure it out. Most didn't. The problem: two Spark jobs "
            "writing to the same prefix with no coordination. No ACID means no rollback, no isolation, "
            "no history. That's why 60% of bank-built lakes became swamps within three years."
        ),
    )


def make_slide_08c_lakehouse(prs):
    """Slide 08c -- Lakehouse Architecture (swimlane diagram)."""
    _make_pattern_swimlane_slide(
        prs,
        title="Lakehouse",
        subtitle="The lake, with ACID bolted on. Iceberg, Delta, Hudi over object storage.",
        highlight_lanes={"sources", "ingestion", "storage", "motion"},
        active_items={
            "sources": ["Core Banking", "CRM", "Card Stream", "Third-Party"],
            "ingestion": ["Batch ETL", "CDC", "Streaming"],
            "storage": ["Lakehouse (Iceberg)", "ACID", "Schema Evolution", "Time Travel", "Partition Pruning"],
            "motion": ["Real-time scoring", "Streaming aggregates"],
        },
        page=10,
        notes=(
            "Same cheap storage, but now with Iceberg on top. ACID means two writers can't corrupt each "
            "other. Schema evolution means adding a column isn't a two-week project. Time travel means the "
            'regulator asks "what did it look like on March 15th?" and you answer in one query instead of '
            "a three-week data recovery project. This is watsonx.data's default format. This is the pitch."
        ),
    )


def make_slide_08d_mesh(prs):
    """Slide 08d -- Data Mesh Architecture (swimlane diagram)."""
    _make_pattern_swimlane_slide(
        prs,
        title="Data Mesh",
        subtitle="Domain-owned data products. The org chart, expressed in storage.",
        highlight_lanes={"sources", "ingestion", "storage", "insight"},
        active_items={
            "sources": ["Retail Domain", "Cards Domain", "Wealth Domain", "Risk Domain"],
            "ingestion": ["Per-domain pipelines", "Data contracts", "Schema registry"],
            "storage": ["Domain Lakehouse (Cards)", "Domain Lakehouse (Retail)", "Domain Lakehouse (Wealth)"],
            "insight": ["Cross-domain analytics", "Data product discovery", "Federated governance"],
        },
        page=11,
        notes=(
            "Mesh is not a technology \u2014 it's an org chart expressed in storage. Each domain owns its "
            "lakehouse and publishes data products with contracts. The platform team provides the runway. "
            'The question isn\'t "should we do mesh?" \u2014 it\'s "does our org have the product thinking '
            "maturity to sustain it?\" Most don't. The ones that do usually have 500+ data engineers and "
            "a sponsoring exec."
        ),
    )


def make_slide_08e_fabric(prs):
    """Slide 08e -- Data Fabric Architecture (swimlane diagram)."""
    _make_pattern_swimlane_slide(
        prs,
        title="Data Fabric",
        subtitle="The meta-pattern. AI-driven governance spanning all other patterns.",
        highlight_lanes={"sources", "ingestion", "storage", "motion", "insight"},
        active_items={
            "sources": ["System of Record", "Application Data", "Third-Party", "Content Services"],
            "ingestion": ["Automated discovery", "AI-classified ingestion", "Policy-aware routing"],
            "storage": ["Warehouse", "Lakehouse", "Lake", "Vector / Search"],
            "motion": ["Automated lineage", "Quality monitoring", "Schema drift detection"],
            "insight": ["Knowledge Catalog", "Self-serve discovery", "Federated governance"],
        },
        page=12,
        notes=(
            "Fabric is the meta-pattern. It doesn't replace the others \u2014 it automates the governance "
            "layer across all of them. The metadata layer knows where everything is, who owns it, how "
            "fresh it is, and what policies apply. AI-assisted classification, automated lineage, self-serve "
            "discovery. This is where IBM's Knowledge Catalog plays. When the customer has a heterogeneous "
            "estate \u2014 Db2, Oracle, Iceberg, and a folder of PDFs \u2014 fabric is the answer. When "
            "they're greenfield single-cloud, it's overkill."
        ),
    )


def make_slide_09_cameo_warehouse(prs):
    """Slide 09 -- Notebook Cameo #1: Warehouse."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 1.3 \u2014 PATTERN DECODER RING (DEMO)",
        "Let\u2019s see one running.",
        page=13,
        chapter="Block 1 \u2014 Foundations",
    )

    add_textbox(
        slide,
        MARGIN_L,
        Inches(1.9),
        CONTENT_W,
        Inches(0.4),
        "NOTEBOOK CAMEO #1",
        font_name=FONT_MONO,
        size=PT_SMALL,
        color=ACCENT,
        bold=True,
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(2.5),
        Inches(10),
        Inches(3.0),
        [
            "Open notebooks/01-warehouse.ipynb",
            "Run one canonical query against the synthetic BFSI star schema",
            "Show the output: branch transaction volume for Q3 2024",
            "Then try Q2 \u2014 finding customers whose policy documents reference an AML procedure",
            "The warehouse stores the PDF as a BLOB. It cannot read it.",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "We have notebooks. They're real. They run. This is the warehouse query \u2014 branch "
            "transaction volume for Q3 2024. Clean SQL, fast results, no surprises. This is what the "
            "warehouse does well. Now watch what happens when we try Q2 \u2014 finding customers whose "
            "policy documents reference an AML procedure. The warehouse stores the PDF as a BLOB. "
            "It cannot read it."
        ),
    )


def make_slide_10_cameo_virtualization(prs):
    """Slide 10 -- Notebook Cameo #2: Virtualization."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 1.3 \u2014 PATTERN DECODER RING (DEMO)",
        "No data moved.",
        page=14,
        chapter="Block 1 \u2014 Foundations",
    )

    add_textbox(
        slide,
        MARGIN_L,
        Inches(1.9),
        CONTENT_W,
        Inches(0.4),
        "NOTEBOOK CAMEO #2",
        font_name=FONT_MONO,
        size=PT_SMALL,
        color=ACCENT,
        bold=True,
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(2.5),
        Inches(10),
        Inches(3.0),
        [
            "Open notebooks/04-virtualization.ipynb",
            "Run one federated query joining two sources without moving data",
            "Joining Postgres and Iceberg data through Trino \u2014 no ETL, no copies, one SQL query",
            "Read the EXPLAIN plan \u2014 that\u2019s where the truth lives",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            'This is the most misunderstood pattern. Everyone nods at "federated queries" in the reference '
            "architecture. Nobody believes it actually works until they see it. Watch: we're joining Postgres "
            "and Iceberg data through Trino. No ETL. No copies. One SQL query. Now \u2014 before you sell "
            "this to every customer \u2014 read the EXPLAIN plan. That's where the truth lives."
        ),
    )


def make_slide_11_mdm(prs):
    """Slide 11 -- MDM: The Unsexy Foundation."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 1.4 \u2014 MASTER DATA MANAGEMENT",
        "The unsexy foundation.",
        page=15,
        chapter="Block 1 \u2014 Foundations",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.0),
        [
            "Why MDM became cool again: agents need authoritative entities to reason over",
            "Four MDM patterns: registry, consolidation, coexistence, centralized",
            "IBM MDM / InfoSphere MDM / Match 360 positioning",
            "Graph + MDM: the entity resolution play (foreshadows Block 2 Graph RAG)",
            "Every Canadian bank has a 10-year MDM program that\u2019s \u201calmost done.\u201d Agentic AI just made it urgent again.",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            'Every Canadian bank has a 10-year MDM program that\'s "almost done." Agentic AI just made it '
            'urgent again. When your agent needs to answer "who is this customer?" across four systems \u2014 '
            "core banking, CRM, AML, branch records \u2014 it needs entity resolution. That's MDM. Not "
            "glamorous. Absolutely foundational. If your customer's MDM is broken, don't sell them RAG. "
            "Sell them MDM first."
        ),
    )


def make_slide_12_ibm_stack(prs):
    """Slide 12 -- IBM Stack: Storage + Access (table)."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 1.5 \u2014 IBM STACK MAPPED TO THE REFERENCE ARCHITECTURE",
        "Where the storage and access patterns live.",
        page=16,
        chapter="Block 1 \u2014 Foundations",
    )

    data = [
        ["LAYER", "IBM PRODUCTS"],
        ["Storage", "watsonx.data, Db2, Db2 Warehouse, Informix, EDB Postgres"],
        ["Compute", "Presto via watsonx.data, Spark, DuckDB"],
        ["Catalog", "Iceberg REST catalog, Knowledge Catalog"],
        ["Federation", "Data Virtualization, Presto connectors"],
        ["Open formats", "Iceberg, Delta Lake, Milvus"],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(1.8),
        CONTENT_W,
        data,
        col_widths=[3.0, 8.8],
        accent_col=0,
    )

    set_notes(
        slide,
        (
            "When you're positioning this with a customer, point at the Analytical Data Management & Storage "
            "lane. watsonx.data is the anchor. If they're on z/OS, it's Db2 for z/OS. If they want modern "
            "open formats, it's Iceberg on watsonx.data with Presto. The federation story is Data "
            "Virtualization or Presto connectors joining sources without copying."
        ),
    )


def make_slide_13_annotated1(prs):
    """Slide 13 -- Annotated Diagram #1: Storage + Access."""
    slide = add_slide(prs)

    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.3),
        CONTENT_W,
        Inches(0.5),
        "Annotated Diagram #1 \u2014 Storage + Access",
        font_name=FONT_DISPLAY,
        size=PT_SUBTITLE,
        color=INK,
        bold=True,
    )
    add_diagram_or_placeholder(
        slide, "refarch-block1.png", "Annotated Ref-Arch: Storage + Access lanes lit", top=Inches(1.0)
    )
    add_footer(slide, "Block 1 \u2014 Foundations", 17)

    set_notes(
        slide,
        (
            "Here's the reference architecture again, but now with only the storage and access lanes lit. "
            "This is what Block 1 was about. watsonx.data sits in the big blue box in the center. The "
            "connectors reach left to the sources and right to the query engines. The governance band at "
            "the bottom \u2014 we'll light that up in Block 3."
        ),
    )


def make_slide_14_transition_quote(prs):
    """Slide 14 -- Transition Quote."""
    slide = add_slide(prs)
    add_big_quote(
        slide,
        "The warehouse works. The lake stores everything. The lakehouse adds ACID. "
        "Virtualization federates. Mesh organizes. None of them can read a PDF.",
        "",
        page=18,
        chapter="Block 1 \u2014 Foundations",
    )
    set_notes(
        slide,
        (
            "That's the setup. Five patterns, each solving a real problem, each with a clear failure mode. "
            "Now let's talk about what AI changed \u2014 and what the reference architecture needs that "
            "isn't on the diagram yet."
        ),
    )


# ── Block 2: AI-Era Architecture ──────────────────────────────────────────────


def make_slide_15_divider_ai(prs):
    """Slide 15 -- Section Divider: AI-Era Architecture."""
    slide = add_slide(prs)
    add_section_divider(
        slide,
        "BLOCK 2 OF 5",
        "AI-Era Architecture.",
        "What changed when AI ate the stack. RAG, Docling, context engineering, and the agent control plane.",
        "60 MIN",
        page=19,
    )
    set_notes(
        slide,
        (
            "If Block 1 was the foundation, Block 2 is why we're all here. AI changed what data "
            "architecture means. New data types, new latency requirements, new governance problems. "
            "Let me show you what changed and what the reference architecture needs that isn't on "
            "the diagram yet."
        ),
    )


def make_slide_16_what_changed(prs):
    """Slide 16 -- What Changed (before/after two-column compare)."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.1 \u2014 WHAT CHANGED",
        "Old vs. new.",
        page=20,
        chapter="Block 2 \u2014 AI-Era",
    )

    data = [
        ["BEFORE (BI ERA)", "AFTER (AI ERA)"],
        ["Scan tables", "Retrieve passages"],
        ["Minutes (batch)", "Milliseconds (agent loops)"],
        ["Structured only", "Unstructured is primary fuel"],
        ["Tables, views", "Embeddings, chunks, traces, prompts, tool-call logs"],
        ["Report-time governance", "Real-time policy enforcement"],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(1.8),
        CONTENT_W,
        data,
        col_widths=[5.9, 5.9],
    )

    set_notes(
        slide,
        (
            "Read patterns flipped. In the BI era, you scanned tables and aggregated. In the AI era, you "
            "retrieve passages and assemble context. Latency budgets collapsed from minutes to milliseconds. "
            "And 80% of enterprise data \u2014 the PDFs, the emails, the contracts \u2014 went from "
            '"unusable for analytics" to "primary fuel for AI." That\'s the unstructured data tax. '
            "Every bank has it. Most haven't paid it yet."
        ),
    )


def make_slide_17_rag_pipeline(prs):
    """Slide 17 -- RAG Pipeline (diagram)."""
    slide = add_slide(prs)

    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.3),
        CONTENT_W,
        Inches(0.3),
        "\u00a7 2.2 \u2014 THE RAG REFERENCE ARCHITECTURE",
        font_name=FONT_MONO,
        size=PT_MONO,
        color=ACCENT,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.65),
        CONTENT_W,
        Inches(0.5),
        "The pipeline that actually ships.",
        font_name=FONT_DISPLAY,
        size=PT_TITLE,
        color=INK,
        bold=True,
    )
    add_diagram_or_placeholder(slide, "rag-pipeline.png", "RAG Pipeline: Docling \u2192 watsonx.ai", top=Inches(1.3))
    add_footer(slide, "Block 2 \u2014 AI-Era", 21)

    set_notes(
        slide,
        (
            "This is the pipeline. Every step can fail. Chunking is where dreams die \u2014 split too "
            "aggressively and you lose context, split too conservatively and you drown the model in noise. "
            "Pure vector retrieval is what consultants sell. Hybrid retrieval \u2014 BM25 plus vector plus "
            "reranker \u2014 is what works. Let me show you each layer."
        ),
    )


def make_slide_18_docling(prs):
    """Slide 18 -- Docling Deep-dive."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.2 \u2014 DOCLING DEEP-DIVE",
        "The ingestion-side fix.",
        page=22,
        chapter="Block 2 \u2014 AI-Era",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.5),
        [
            "Open-source, Apache 2.0, IBM Research Zurich, now in Linux Foundation\u2019s Agentic AI Foundation",
            "Replaces naive OCR with layout-aware extraction (~30x speedup, preserves structure)",
            "DoclingDocument representation: bounding boxes, reading order, structure-aware chunking",
            "Integrates with LangChain, LlamaIndex, spaCy",
            "Granite-Docling-258M VLM under Apache 2.0",
            "Docling OpenShift Operator with Red Hat \u2014 banks named as deployment segment",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            'Docling is the answer to "how do I get PDFs into my RAG pipeline without losing the table '
            'structure?" Naive OCR reads left to right and destroys multi-column layouts. Docling '
            "understands the visual structure \u2014 it knows that the table on page 7 has four columns, "
            "not eight. The OpenShift Operator means your bank can run it on-prem, air-gapped, with no "
            "data leaving the building. That's the pitch."
        ),
    )


def make_slide_19_opensearch(prs):
    """Slide 19 -- OpenSearch Hybrid Retrieval."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.2 \u2014 HYBRID RETRIEVAL",
        "BM25 + vector + reranker. In one engine.",
        page=23,
        chapter="Block 2 \u2014 AI-Era",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.5),
        [
            "Why pure vector search fails: semantic similarity is not semantic correctness",
            "BM25 for exact matches (regulatory terms, policy numbers)",
            "Vector for conceptual similarity",
            "Reranker for precision (cross-encoder on top-k)",
            "OpenSearch does all three in one engine \u2014 no stitching",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "Pure vector RAG is a myth. I've never seen it work in production at a bank. The compliance "
            'officer searches for "PCMLTFA Section 7 obligation for EFTs over $1000" \u2014 that\'s a '
            "keyword search, not a semantic one. You need BM25 for exact terms AND vector for conceptual "
            "similarity AND a reranker to sort the results. OpenSearch does all three. That's why it's "
            "in the reference architecture."
        ),
    )


def make_slide_20_vector_stores(prs):
    """Slide 20 -- Vector Store Decision Tree (table)."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.2 \u2014 CHOOSING A VECTOR STORE",
        "Pick based on ops maturity, not benchmarks.",
        page=24,
        chapter="Block 2 \u2014 AI-Era",
    )

    data = [
        ["STORE", "BEST FOR", "WATCH OUT FOR"],
        ["Milvus", "Scale, GPU-accelerated search", "Ops complexity"],
        ["OpenSearch", "Hybrid (BM25 + vector + reranker)", "Memory footprint"],
        ["Elastic", "Existing ELK stack shops", "License changes"],
        ["Pinecone", "Managed, zero-ops", "Vendor lock-in, no hybrid"],
        ["pgvector", "Already on Postgres", "Not built for scale"],
        ["Chroma", "Prototyping", "Not production-grade"],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(1.8),
        CONTENT_W,
        data,
        col_widths=[2.5, 4.7, 4.6],
        accent_col=0,
    )

    set_notes(
        slide,
        (
            "Stop benchmarking vector stores. Start asking: does my team know how to operate this in "
            "production? Pinecone is zero-ops but vendor-locked. OpenSearch is self-managed but does "
            "hybrid. pgvector is free but not built for scale. Pick based on your ops maturity, not "
            "someone else's benchmark."
        ),
    )


def make_slide_21_context_engineering(prs):
    """Slide 21 -- Beyond RAG: Context Engineering."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.3 \u2014 CONTEXT ENGINEERING AND THE AGENT DATA PLANE",
        "RAG vs. fine-tuning is the wrong question.",
        page=25,
        chapter="Block 2 \u2014 AI-Era",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.5),
        [
            "Context engineering: what goes in the window, in what order, with what provenance",
            "Agent data plane: tool registries, memory stores, trace logs, evaluation sets \u2014 these are data architecture now",
            "Graph RAG and when it beats vector (entity-heavy, multi-hop reasoning)",
            "Open RAG positioning: open standards, open formats, swappable components",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            '"Should we RAG or fine-tune?" is the question every customer asks. The answer is: neither, '
            "both, and it depends. The real question is context engineering \u2014 what goes in the context "
            "window, in what order, with what provenance. That's a data architecture problem. Tool "
            "registries, memory stores, trace logs, eval sets \u2014 these are the new data products. "
            "If you're an architect, this is your job now."
        ),
    )


def make_slide_22_agent_data_plane(prs):
    """Slide 22 -- Agent Data Plane Components (table)."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.3 \u2014 THE AGENT DATA PLANE",
        "These are data architecture now.",
        page=26,
        chapter="Block 2 \u2014 AI-Era",
    )

    data = [
        ["COMPONENT", "WHAT IT STORES", "WHY IT MATTERS"],
        ["Tool registry", "Available tools, permissions, schemas", "Agent needs to know what it can call"],
        ["Memory store", "Conversation history, entity state", "Agents need context across sessions"],
        ["Trace logs", "Every tool call, every generation", "Audit, debugging, compliance"],
        ["Evaluation sets", "Ground-truth Q&A pairs, human ratings", "Only way to know if it works"],
        ["Prompt templates", "Versioned prompts, system messages", "Prompt is code; version it like code"],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(1.8),
        CONTENT_W,
        data,
        col_widths=[2.8, 4.2, 4.8],
        accent_col=0,
    )

    set_notes(
        slide,
        (
            "Five new data products that didn't exist three years ago. Every one of these needs storage, "
            "versioning, access control, and lineage. If you're a data architect, this is your new backlog. "
            "If you're a seller, this is the expansion play after the initial RAG deployment."
        ),
    )


def make_slide_23_open_rag(prs):
    """Slide 23 -- Open RAG Positioning."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.3 \u2014 OPEN RAG",
        "Open standards, open formats, swappable components.",
        page=27,
        chapter="Block 2 \u2014 AI-Era",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.5),
        [
            "Enterprise buyers are demanding component swappability",
            "No lock-in to one embedding model, one vector store, one reranker",
            "Open formats: Parquet for structured, Docling for unstructured, OpenSearch for retrieval",
            "IBM\u2019s position: watsonx.ai + watsonx.data + OpenSearch is the open stack",
            "Contrast with closed alternatives (Azure AI Search + Azure OpenAI, Bedrock + Kendra)",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "Enterprise buyers are done with lock-in. They want to swap the embedding model without "
            "rewriting the pipeline. They want to change the vector store without losing the index. "
            "Open RAG means: Parquet for structured data, Docling for document extraction, OpenSearch "
            "for hybrid retrieval, and any LLM behind the generation step. IBM's play is the open "
            "stack. Position it against the closed alternatives where every component is vendor-locked."
        ),
    )


def make_slide_24_mcp(prs):
    """Slide 24 -- MCP: The Integration Contract."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.4 \u2014 MCP AND THE AGENT CONTROL PLANE",
        "The integration contract between agents and your enterprise.",
        page=28,
        chapter="Block 2 \u2014 AI-Era",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.5),
        [
            "Why MCP matters for data architects: standardized tool interface for agents",
            "The problem it solves: every agent framework has its own tool-calling convention",
            "MCP as the \u201cJDBC of agents\u201d \u2014 one contract, many implementations",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "MCP \u2014 Model Context Protocol \u2014 is what JDBC was for databases. One standard interface "
            "so agents can call tools without caring about the implementation. Before JDBC, every database "
            "had its own wire protocol and every app was hardcoded to one vendor. Before MCP, every agent "
            "framework had its own tool-calling convention. MCP standardizes it. If you're selling data "
            "architecture, MCP is how your customer's data becomes accessible to agents."
        ),
    )


def make_slide_25_context_forge(prs):
    """Slide 25 -- MCP Context Forge: The Missing Layer."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.4 \u2014 CONTEXT FORGE",
        "The missing layer between AI agents and your enterprise.",
        page=29,
        chapter="Block 2 \u2014 AI-Era",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.5),
        [
            "NHI attribution (non-human identity \u2014 which agent, which tool, which user delegated)",
            "Blast radius control (what can this agent touch?)",
            "HITL gate (human-in-the-loop for high-risk actions)",
            "DLP (data loss prevention \u2014 PII doesn\u2019t leave the perimeter)",
            "Immutable audit (every tool call logged, tamper-evident)",
            "Anomaly detection (agent doing something unusual? flag it)",
            "This is a gap in the published reference architecture \u2014 call it out explicitly",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "Full disclosure: I'm a contributor to Context Forge. This is the missing layer in the "
            "reference architecture \u2014 and I say that knowing the published diagram doesn't have it "
            "yet. When your agent calls a tool that queries customer data, who approved it? What's the "
            "blast radius? Where's the audit trail? Context Forge answers these questions. The plugin "
            "model means you can add governance capabilities without rewriting your agent."
        ),
    )


def make_slide_26_notebook_teaser(prs):
    """Slide 26 -- Notebook Teaser #3: Docling Live."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 2.2 \u2014 DOCLING LIVE DEMO",
        "Structure-aware extraction in action.",
        page=30,
        chapter="Block 2 \u2014 AI-Era",
    )

    add_textbox(
        slide,
        MARGIN_L,
        Inches(1.9),
        CONTENT_W,
        Inches(0.4),
        "NOTEBOOK TEASER #3",
        font_name=FONT_MONO,
        size=PT_SMALL,
        color=ACCENT,
        bold=True,
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(2.5),
        Inches(10),
        Inches(3.5),
        [
            "Open notebooks/06-rag-mdm.ipynb to the Docling cell",
            "Parse a real PDF live \u2014 watch the structured output with section headings preserved",
            "This is what goes into the chunking step",
            "Compare with naive 512-character window splitting \u2014 this is the fix",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "Let me show you Docling in action. This is one of our 10 synthetic AML policy PDFs. Watch "
            "the output \u2014 it preserves the section headings, the table structure, the reading order. "
            "This is what goes into the chunking step. If you've ever seen a RAG pipeline that chops "
            "documents into 512-character windows regardless of structure \u2014 this is the fix."
        ),
    )


def make_slide_27_annotated2(prs):
    """Slide 27 -- Annotated Diagram #2: AI + Ingestion."""
    slide = add_slide(prs)

    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.3),
        CONTENT_W,
        Inches(0.5),
        "Annotated Diagram #2 \u2014 AI + Ingestion",
        font_name=FONT_DISPLAY,
        size=PT_SUBTITLE,
        color=INK,
        bold=True,
    )
    add_diagram_or_placeholder(
        slide, "refarch-block2.png", "Annotated Ref-Arch: AI + Ingestion lanes lit", top=Inches(1.0)
    )
    add_footer(slide, "Block 2 \u2014 AI-Era", 31)

    set_notes(
        slide,
        (
            "Here's the diagram again. Now we've lit up the AI and analytics lanes \u2014 Watson Studio, "
            "watsonx.ai, the ML pipelines. Notice two annotations I've added. First: Docling, sitting in "
            "the Ingestion & Integration lane \u2014 it's how unstructured documents enter the pipeline. "
            "Second: Context Forge, below the governance band \u2014 it's the agent control plane that "
            "the published diagram doesn't include yet. You're ahead of the published material. Use that "
            "in customer conversations."
        ),
    )


def make_slide_28_transition_quote2(prs):
    """Slide 28 -- Transition Quote."""
    slide = add_slide(prs)
    add_big_quote(
        slide,
        "Every bank has data governance. Almost none have AI governance. Zero have agent governance. Those are three different problems.",
        "",
        page=32,
        chapter="Block 2 \u2014 AI-Era",
    )
    set_notes(
        slide,
        (
            "Let that sink in. Three different problems, three different toolchains, three different teams. "
            "After the break, we're going to light up the bottom three bands of the reference architecture "
            'and talk about what "governed" actually means in Canada.'
        ),
    )


# ── Block 3: Governance ───────────────────────────────────────────────────────


def make_slide_29_divider_governance(prs):
    """Slide 29 -- Section Divider: Governance."""
    slide = add_slide(prs)
    add_section_divider(
        slide,
        "BLOCK 3 OF 5",
        "Governance.",
        "Data governance. AI governance. Agent governance. Three problems, not one.",
        "45 MIN",
        page=33,
    )
    set_notes(
        slide,
        (
            "Welcome back. This block is the one that makes the compliance team lean forward. We're going "
            'to talk about what "regulated" actually means, why every reference architecture you see '
            "online is wrong for a Canadian bank, and why agents just made the threat model worse."
        ),
    )


def make_slide_30_governance_triad(prs):
    """Slide 30 -- Three Governance Problems (diagram)."""
    slide = add_slide(prs)

    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.3),
        CONTENT_W,
        Inches(0.3),
        "\u00a7 3.1 \u2014 THREE GOVERNANCE PROBLEMS",
        font_name=FONT_MONO,
        size=PT_MONO,
        color=ACCENT,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.65),
        CONTENT_W,
        Inches(0.5),
        "Treating these as one thing is how enterprises get burned.",
        font_name=FONT_DISPLAY,
        size=PT_TITLE,
        color=INK,
        bold=True,
    )
    add_diagram_or_placeholder(
        slide, "governance-triad.png", "Governance Triad: Data / AI / Agent maturity", top=Inches(1.3)
    )
    add_footer(slide, "Block 3 \u2014 Governance", 34)

    set_notes(
        slide,
        (
            "Data governance is table-stakes \u2014 every bank has Knowledge Catalog, lineage, quality "
            "rules. AI governance is emerging \u2014 model risk management, OSFI E-23, bias detection, "
            "drift monitoring. Agent governance is brand new \u2014 who authorized this agent to call this "
            "tool? What's the blast radius? Where's the kill switch? Treating these as one problem is how "
            "you end up with a governance framework that covers data lineage perfectly but has no idea "
            "an agent just exfiltrated PII via a tool call."
        ),
    )


def make_slide_31_regulated(prs):
    """Slide 31 -- What Regulated Actually Means."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 3.2 \u2014 REGULATED REALITY",
        "Why every reference architecture you see online is wrong for a Canadian bank.",
        page=35,
        chapter="Block 3 \u2014 Governance",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.0),
        [
            "OSFI B-13 (technology and cyber risk)",
            "PIPEDA (privacy)",
            "E-23 (model risk management)",
            "SR 11-7 (US model risk, but Canadian banks follow it too)",
            "Lineage-to-audit-trail requirements",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "When I show a reference architecture at a US tech conference, everyone nods. When I show it "
            'at a Canadian bank, the CISO asks "where\'s the data residency boundary?" and the model risk '
            'officer asks "where\'s the E-23 mapping?" and the privacy officer asks "where\'s the PIPEDA '
            'consent flow?" Every reference architecture you see online is wrong for a Canadian bank. '
            "Today we fix that."
        ),
    )


def make_slide_32_sovereignty(prs):
    """Slide 32 -- Sovereignty vs. Residency vs. Operational Sovereignty."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 3.2 \u2014 SOVEREIGNTY",
        "Three different things, frequently confused.",
        page=36,
        chapter="Block 3 \u2014 Governance",
    )

    data = [
        ["TERM", "DEFINITION"],
        ["Data Residency", "Data physically stays in Canada"],
        ["Data Sovereignty", "Canadian law governs the data"],
        ["Operational Sovereignty", "Canadian operators control the infrastructure"],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(1.8),
        Inches(10),
        data,
        col_widths=[3.5, 6.5],
        accent_col=0,
    )

    # Pull quote below table
    add_textbox(
        slide,
        MARGIN_L,
        Inches(4.2),
        Inches(10),
        Inches(0.6),
        "\u201cSovereignty, not solitude.\u201d \u2014 Sovereignty doesn\u2019t mean isolation.",
        font_name=FONT_DISPLAY,
        size=PT_BODY,
        color=INK2,
        italic=True,
    )

    set_notes(
        slide,
        (
            "Three different things, frequently confused. A Canadian bank can use AWS in Montreal "
            "(residency) but the data is still subject to the US CLOUD Act (sovereignty gap). "
            "Operational sovereignty means a Canadian team runs the infrastructure \u2014 not an "
            'offshore NOC. "Sovereignty, not solitude" \u2014 you can use public cloud, but you need '
            "to know which of these three you're actually getting."
        ),
    )


def make_slide_33_observability(prs):
    """Slide 33 -- Data Observability."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 3.3 \u2014 DATA OBSERVABILITY",
        "The triad nobody has all three of.",
        page=37,
        chapter="Block 3 \u2014 Governance",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.5),
        [
            "Data observability: quality, freshness, schema drift, pipeline SLOs",
            "Model observability: bias, drift, performance degradation",
            "Agent trace observability: tool call logs, latency, error rates, anomaly detection",
            "Tools: Monte Carlo, Bigeye, IBM Databand, OpenLineage",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "The triad: data observability, model observability, agent trace observability. Three layers, "
            "frequently confused, rarely all instrumented. Most banks have data observability (at least for "
            "their warehouse pipelines). Some have model observability (watsonx.governance does this). "
            "Almost none have agent trace observability \u2014 logging every tool call, every context "
            "assembly, every generation. That's the gap."
        ),
    )


def make_slide_34_10_domain(prs):
    """Slide 34 -- 10-Domain Agentic Governance Framework."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 3.4 \u2014 GOVERNANCE AS ARCHITECTURE",
        "75 sub-capabilities, 10 domains. Overview only \u2014 we\u2019ll go deep another day.",
        page=38,
        chapter="Block 3 \u2014 Governance",
    )

    # Display the 10 domains as a grid (2 rows x 5 cols)
    domains = [
        "Identity",
        "Authorization",
        "Audit",
        "Evaluation",
        "Safety",
        "Privacy",
        "Quality",
        "Observability",
        "Lifecycle",
        "Compliance",
    ]
    box_w = Inches(2.1)
    box_h = Inches(0.8)
    gap_x = Inches(0.25)
    gap_y = Inches(0.3)
    start_x = MARGIN_L + Inches(0.2)
    start_y = Inches(2.2)

    for i, domain in enumerate(domains):
        row, col = divmod(i, 5)
        x = start_x + col * (box_w + gap_x)
        y = start_y + row * (box_h + gap_y)
        add_rect(slide, x, y, box_w, box_h, fill=ACCENT, line_color=None)
        add_textbox(
            slide,
            x,
            y + Inches(0.15),
            box_w,
            box_h - Inches(0.3),
            domain,
            font_name=FONT_MONO,
            size=PT_BODY,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    set_notes(
        slide,
        (
            "The full framework has 75 sub-capabilities across 10 domains. We don't have time to go "
            'through all of them today, but I want you to know it exists. When a customer asks "how do '
            "we govern our agents?\" \u2014 this is the answer. It maps to a real bank's 80-domain "
            "questionnaire. Context Forge implements it as plugins."
        ),
    )


def make_slide_35_bank_questionnaire(prs):
    """Slide 35 -- Bank 80-Domain Questionnaire Mapping."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 3.4 \u2014 BANK QUESTIONNAIRE MAPPING",
        "How RBC, CIBC, and Scotia approach this differently.",
        page=39,
        chapter="Block 3 \u2014 Governance",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.0),
        [
            "Real bank governance questionnaires have 60\u201380 domains",
            "The 10-domain agentic governance framework maps to the top 15\u201320",
            "Anonymized patterns: one bank leads with privacy, another with model risk, a third with operational risk",
            "\u201cSovereignty, not solitude\u201d applied to governance: federate where you can, centralize where you must",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "I can't name the banks. But I can tell you the pattern. One Big Five bank builds governance "
            "from privacy outward. Another starts from model risk (E-23 compliance). A third starts from "
            "operational risk. They all end up in the same place \u2014 a framework that covers data, AI, "
            "and agent governance \u2014 but they get there from different starting points. When you're in "
            "front of the customer, ask which door they came through. That tells you how to position."
        ),
    )


def make_slide_36_factsheets(prs):
    """Slide 36 -- AI Factsheets and Model Inventory."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 3.4 \u2014 WATSONX.GOVERNANCE",
        "Where HITL and BPM integration fits.",
        page=40,
        chapter="Block 3 \u2014 Governance",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.0),
        [
            "watsonx.governance: model inventory, AI Factsheets, bias/drift monitoring",
            "Orchestrate: where human-in-the-loop and BPM integration happens",
            "Plugin-based governance: how Context Forge extends the framework",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "watsonx.governance is where model risk management lives. AI Factsheets document what the "
            "model does, who trained it, what data it saw, and what its known failure modes are. "
            "Orchestrate handles the human-in-the-loop workflows \u2014 when the model says \"I'm not "
            'sure," Orchestrate routes it to a human. Context Forge extends this to agents \u2014 same '
            "idea, but for autonomous tool calls instead of model predictions."
        ),
    )


def make_slide_37_threat_model(prs):
    """Slide 37 -- Why Agents Change the Threat Model."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 3.5 \u2014 AGENT THREAT MODEL",
        "Autonomous tool use is a new attack surface.",
        page=41,
        chapter="Block 3 \u2014 Governance",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.5),
        [
            "Lateral movement via tool chains",
            "Prompt injection as exfiltration vector",
            "Control plane requirements: identity, authorization, audit, kill switch, eval harness",
            "Where the partner ecosystem (AWS, Google, Microsoft, JPMC, Palo Alto) is converging",
            "Where IBM is differentiated \u2014 and where we\u2019re catching up",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "When an agent calls a tool, it's executing code with privileges. If the agent is compromised "
            "via prompt injection, the attacker inherits those privileges. That's lateral movement via tool "
            "chain \u2014 the agent becomes the attack vector. This is why the control plane needs identity "
            "(who is this agent acting for?), authorization (what tools can it call?), audit (what did it "
            "actually call?), and a kill switch (stop everything, now). The named-partner ecosystem is "
            "converging on these requirements. We're not ahead of everyone, but we're not behind either."
        ),
    )


def make_slide_38_annotated3(prs):
    """Slide 38 -- Annotated Diagram #3: Governance + Security + Deploy."""
    slide = add_slide(prs)

    add_textbox(
        slide,
        MARGIN_L,
        Inches(0.3),
        CONTENT_W,
        Inches(0.5),
        "Annotated Diagram #3 \u2014 Governance + Security + Deploy",
        font_name=FONT_DISPLAY,
        size=PT_SUBTITLE,
        color=INK,
        bold=True,
    )
    add_diagram_or_placeholder(
        slide, "refarch-block3.png", "Annotated Ref-Arch: Governance + Security + Deploy bands lit", top=Inches(1.0)
    )
    add_footer(slide, "Block 3 \u2014 Governance", 42)

    set_notes(
        slide,
        (
            "Last time with the diagram. Now we've lit up the bottom three bands \u2014 governance, "
            "security, and deployment. This is where IBM's differentiation is strongest. Knowledge "
            "Catalog, Data Lineage, MDM/Match 360, AI Factsheets, Guardium, OpenShift. When the CISO "
            'asks "who controls this?" \u2014 point here.'
        ),
    )


def make_slide_39_ibm_strengths_gaps(prs):
    """Slide 39 -- Where IBM's POV Is Differentiated (strengths/gaps)."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 3.5 \u2014 IBM DIFFERENTIATION",
        "Where we lead and where we\u2019re catching up.",
        page=43,
        chapter="Block 3 \u2014 Governance",
    )

    data = [
        ["STRENGTH", "GAP"],
        ["Knowledge Catalog + lineage", "Agent control plane (Context Forge fills this)"],
        ["watsonx.governance + AI Factsheets", "Agent trace observability"],
        ["OpenShift everywhere", "Competitive vector store (Milvus vs. managed options)"],
        ["MDM/Match 360", "Open RAG standardization"],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(1.8),
        CONTENT_W,
        data,
        col_widths=[5.9, 5.9],
    )

    set_notes(
        slide,
        (
            "I'm going to be honest about where we're strong and where we're catching up. Knowledge "
            "Catalog is best-in-class for data governance. watsonx.governance is strong for model risk. "
            "MDM/Match 360 is the entity resolution play nobody else has. Where we're catching up: agent "
            "control plane (Context Forge is the answer, but it's early), agent trace observability, and "
            "the competitive vector store landscape. When you're in a customer conversation, lead with "
            "governance \u2014 that's where we win."
        ),
    )


# ── Block 4: Hands-on ─────────────────────────────────────────────────────────


def make_slide_40_divider_handson(prs):
    """Slide 40 -- Section Divider: Hands-on."""
    slide = add_slide(prs)
    add_section_divider(
        slide,
        "BLOCK 4 OF 5",
        "Hands-on.",
        "Two notebooks. One critique. Your turn.",
        "60 MIN",
        page=44,
    )
    set_notes(
        slide,
        (
            "Your turn. We're going to run two notebooks live \u2014 the lakehouse and the AI-era "
            "capstone. Then you're going to critique an architecture that someone actually shipped to a "
            "healthcare customer. Buckle up."
        ),
    )


def make_slide_41_setup(prs):
    """Slide 41 -- Setup: The BFSI Scenario."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 4.1 \u2014 SETUP AND ARCHITECTURE MAP",
        "Maple Trust Bank: an AML policy Q&A system with full lineage and audit.",
        page=45,
        chapter="Block 4 \u2014 Hands-on",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(3.5),
        [
            "Recap the reference architecture one more time",
            "You\u2019ve now seen all five swimlane groups lit up",
            "Now we exercise them in code",
            "The scenario: Maple Trust Bank needs an AML policy Q&A system with full lineage and audit",
            "Two notebooks, two patterns, one story",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "One more look at the diagram. You've now seen all five swimlane groups lit up. Now we're "
            "going to exercise them in code. The scenario: Maple Trust Bank needs an AML policy Q&A "
            "system with full lineage and audit. Two notebooks, two patterns, one story."
        ),
    )


def make_slide_42_notebook_map(prs):
    """Slide 42 -- Notebook Map (table)."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 4.1 \u2014 NOTEBOOK MAP",
        "Which notebook lights which swimlane.",
        page=46,
        chapter="Block 4 \u2014 Hands-on",
    )

    data = [
        ["NOTEBOOK", "SWIMLANE", "MODE"],
        ["01 Warehouse", "Analytical Storage", "Take-home (cameo in Block 1)"],
        ["02 Data Lake", "Raw Storage", "Take-home"],
        ["03 Lakehouse", "Analytical Storage + Catalog", "Live now"],
        ["04 Virtualization", "Data Access", "Take-home (cameo in Block 1)"],
        ["05 Data Mesh", "Data Products + Governance", "Take-home"],
        ["06 MDM + RAG", "AI + Governance + Storage", "Live now"],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(1.8),
        CONTENT_W,
        data,
        col_widths=[3.0, 4.4, 4.4],
        accent_col=0,
    )

    set_notes(
        slide,
        (
            "Here's the map. Each notebook lights up different swimlanes. We're going to run Notebook 3 "
            "(Lakehouse) and Notebook 6 (MDM + RAG) live. The other four are take-home \u2014 run them "
            "tonight, break them deliberately, build intuition."
        ),
    )


def make_slide_43_live_nb3(prs):
    """Slide 43 -- Live: Notebook 3 (Lakehouse)."""
    slide = add_slide(prs)

    add_textbox(
        slide,
        MARGIN_L,
        Inches(2.2),
        CONTENT_W,
        Inches(1.0),
        "Live: Notebook 3",
        font_name=FONT_DISPLAY,
        size=36,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(3.4),
        CONTENT_W,
        Inches(0.6),
        "Lakehouse Pattern \u2014 15 minutes",
        font_name=FONT_DISPLAY,
        size=PT_SUBTITLE,
        color=INK2,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(4.4),
        CONTENT_W,
        Inches(0.5),
        "notebooks/03-lakehouse.ipynb",
        font_name=FONT_MONO,
        size=PT_BODY,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Block 4 \u2014 Hands-on", 47)

    set_notes(
        slide,
        (
            "Run cells top to bottom. Narrate the Iceberg schema evolution and time travel cells. Pause "
            "for questions after the partition pruning demo. If a cell fails, skip to pre-recorded output. "
            "(Full choreography will be in facilitator-guide.md \u2014 see Issue #16.)"
        ),
    )


def make_slide_44_live_nb6(prs):
    """Slide 44 -- Live: Notebook 6 (MDM + RAG)."""
    slide = add_slide(prs)

    add_textbox(
        slide,
        MARGIN_L,
        Inches(2.2),
        CONTENT_W,
        Inches(1.0),
        "Live: Notebook 6",
        font_name=FONT_DISPLAY,
        size=36,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(3.4),
        CONTENT_W,
        Inches(0.6),
        "MDM + RAG \u2014 25 minutes",
        font_name=FONT_DISPLAY,
        size=PT_SUBTITLE,
        color=INK2,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        MARGIN_L,
        Inches(4.4),
        CONTENT_W,
        Inches(0.5),
        "notebooks/06-rag-mdm.ipynb",
        font_name=FONT_MONO,
        size=PT_BODY,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Block 4 \u2014 Hands-on", 48)

    set_notes(
        slide,
        (
            "The payoff of the whole day. Run the Docling PDF parsing cell first (takes 30-60s). Then "
            "the OpenSearch indexing. Then Q2 \u2014 the question every previous notebook deferred. Let "
            "the room see it work. Then show Break 2 (the eval scoring) to prove we're honest about "
            "the limits. (Full choreography will be in facilitator-guide.md \u2014 see Issue #16.)"
        ),
    )


def make_slide_45_critique(prs):
    """Slide 45 -- Architecture Critique."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "\u00a7 4.4 \u2014 ARCHITECTURE CRITIQUE",
        "Find the failures.",
        page=49,
        chapter="Block 4 \u2014 Hands-on",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.5),
        [
            "Distribute handouts/critique.md \u2014 HealthFirst Insights scenario",
            "Simple RAG sold as \u201cagentic,\u201d tightly coupled .NET/C#, vibes-based eval",
            "PHI sent to a hyperscaler API with no guardrail, multi-tenancy \u201cplanned\u201d at DB level",
            "Small groups (3\u20134 people), 10 min to find failures, 5 min readout",
            "Severity: Critical / High / Medium    Category: Data / AI / Governance / Security / Ops",
        ],
        size=PT_BODY,
        color=INK,
    )

    set_notes(
        slide,
        (
            "Names and details are anonymized. The pattern is real. Welcome to enterprise AI. You have "
            "10 minutes with your group. Use the worksheet. Severity: critical, high, medium. Category: "
            "data, AI, governance, security, ops. Go."
        ),
    )


# ── Block 5: Close & Q&A ──────────────────────────────────────────────────────


def make_slide_46_takeaway(prs):
    """Slide 46 -- The One-Slide Takeaway (Big Quote)."""
    slide = add_slide(prs)
    add_big_quote(
        slide,
        "Data architecture in 2026 is the integration contract between your enterprise and your agents. Get it wrong and nothing else matters.",
        "",
        page=50,
        chapter="Block 5 \u2014 Close",
    )
    set_notes(
        slide,
        (
            'One slide. One sentence. This is the thing you take back to your customer. Not "buy watsonx" '
            "\u2014 that's the how. The what is: your data architecture IS the integration contract between "
            "your enterprise and your agents. If the contract is broken, nothing else works. If the contract "
            "is solid, everything else follows."
        ),
    )


def make_slide_47_monday_actions(prs):
    """Slide 47 -- Three Monday Morning Actions."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "BLOCK 5 \u2014 CLOSE",
        "Three things to do Monday morning.",
        page=51,
        chapter="Block 5 \u2014 Close",
    )

    data = [
        ["ROLE", "ACTION"],
        ["Sellers", "Stop selling \u201cAI.\u201d Start selling the data plane that makes AI work."],
        ["Architects", "Audit your context engineering and your control plane, not just your RAG pipeline."],
        ["SSRs", "Every incident is now a governance incident."],
    ]
    add_table(
        slide,
        MARGIN_L,
        Inches(2.2),
        Inches(10),
        data,
        col_widths=[2.5, 7.5],
        accent_col=0,
    )

    set_notes(
        slide,
        (
            'Three things. Do them Monday. Sellers: the next time you\'re about to say "AI" in a pitch, '
            'replace it with "data plane." Architects: audit your context engineering, not just your RAG '
            "pipeline. SSRs: every outage is a governance question now \u2014 which swimlane broke, and "
            "who owns it?"
        ),
    )


def make_slide_48_pointers(prs):
    """Slide 48 -- Pointers (links list)."""
    slide = add_slide(prs)
    add_field_guide_header(
        slide,
        "BLOCK 5 \u2014 CLOSE",
        "Pointers.",
        page=52,
        chapter="Block 5 \u2014 Close",
    )

    add_bullet_list(
        slide,
        MARGIN_L,
        Inches(1.9),
        Inches(10),
        Inches(4.0),
        [
            "MCP Context Forge: github.com/IBM/mcp-context-forge",
            "factor10.ai \u2014 10 Principles of Enterprise AI",
            "Docling: github.com/DS4SD/docling",
            "OpenSearch hybrid retrieval guide",
            "Internal IBM governance framework doc",
        ],
        size=PT_BODY,
        color=ACCENT,
    )

    set_notes(
        slide,
        (
            "All of these are live links. The notebooks are on GitHub. Context Forge is open source. "
            "Docling is open source. factor10.ai is the 10 Principles of Enterprise AI. Bookmark them, "
            "share them with your customer's architect. These are the resources that make you credible."
        ),
    )


def make_slide_49_closing(prs):
    """Slide 49 -- Closing."""
    slide = add_slide(prs)

    # Accent blue header bar at top
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.4), fill=ACCENT)

    # Top label
    add_textbox(
        slide,
        MARGIN_L,
        Inches(1.8),
        CONTENT_W,
        Inches(0.4),
        "END OF VOLUME ONE",
        font_name=FONT_MONO,
        size=PT_MONO,
        color=SLATE,
    )

    # Title
    add_textbox(
        slide,
        MARGIN_L,
        Inches(2.5),
        CONTENT_W,
        Inches(1.0),
        "Six patterns. One conversation.",
        font_name=FONT_DISPLAY,
        size=36,
        color=INK,
        bold=True,
    )

    # Subtitle (italic)
    add_textbox(
        slide,
        MARGIN_L,
        Inches(3.6),
        Inches(10),
        Inches(1.5),
        "Walk into the next customer with the pattern decoder ring. "
        "Listen for which row their problem sits on. Propose the pattern, not the product.",
        font_name=FONT_DISPLAY,
        size=PT_SUBTITLE,
        color=INK2,
        italic=True,
    )

    # Bottom rule
    add_rule(slide, MARGIN_L, Inches(5.8), CONTENT_W, color=RULE, weight=1.0)

    # Bottom bar
    add_textbox(
        slide,
        MARGIN_L,
        Inches(6.0),
        CONTENT_W,
        Inches(0.3),
        "QUESTIONS?    |    MANAV.GUPTA@IBM.COM    |    NOTEBOOKS \u00b7 GITHUB.COM/...",
        font_name=FONT_MONO,
        size=PT_MONO,
        color=SLATE,
    )

    add_footer(slide, "Close", 53)

    set_notes(
        slide,
        (
            "Thank you. The notebooks are on GitHub. The pattern decoder ring is your cheat sheet. When "
            'the customer says "we need to query across three systems without moving data," you know '
            'the answer is virtualization. When they say "we need ACID on the lake," it\'s lakehouse. '
            'When they say "our domains need autonomy," it\'s mesh \u2014 but only if they have the org '
            "maturity. Propose the pattern. Then back it with the product."
        ),
    )


# ── Main: assemble all slides ─────────────────────────────────────────────────

SLIDE_BUILDERS = [
    # Opening (01-04)
    make_slide_01_cover,  # 1
    make_slide_02_hook,  # 2
    make_slide_03_refarch,  # 3
    make_slide_03b_refarch_overview,  # 3b — color-coded swimlane overview
    make_slide_03c_refarch_products,  # 3c — IBM product mapping
    make_slide_04_roles,  # 4
    # Block 1 — Foundations (05-14 + 08b-08e)
    make_slide_05_divider_foundations,  # 5
    make_slide_06_reinvention,  # 6
    make_slide_07_primitives,  # 7
    make_slide_08_decoder_ring,  # 8
    make_slide_08b_lake,  # 9  (08b)
    make_slide_08c_lakehouse,  # 10 (08c)
    make_slide_08d_mesh,  # 11 (08d)
    make_slide_08e_fabric,  # 12 (08e)
    make_slide_09_cameo_warehouse,  # 13
    make_slide_10_cameo_virtualization,  # 14
    make_slide_11_mdm,  # 15
    make_slide_12_ibm_stack,  # 16
    make_slide_13_annotated1,  # 17
    make_slide_14_transition_quote,  # 18
    # Block 2 — AI-Era (15-28)
    make_slide_15_divider_ai,  # 19
    make_slide_16_what_changed,  # 20
    make_slide_17_rag_pipeline,  # 21
    make_slide_18_docling,  # 22
    make_slide_19_opensearch,  # 23
    make_slide_20_vector_stores,  # 24
    make_slide_21_context_engineering,  # 25
    make_slide_22_agent_data_plane,  # 26
    make_slide_23_open_rag,  # 27
    make_slide_24_mcp,  # 28
    make_slide_25_context_forge,  # 29
    make_slide_26_notebook_teaser,  # 30
    make_slide_27_annotated2,  # 31
    make_slide_28_transition_quote2,  # 32
    # Block 3 — Governance (29-39)
    make_slide_29_divider_governance,  # 33
    make_slide_30_governance_triad,  # 34
    make_slide_31_regulated,  # 35
    make_slide_32_sovereignty,  # 36
    make_slide_33_observability,  # 37
    make_slide_34_10_domain,  # 38
    make_slide_35_bank_questionnaire,  # 39
    make_slide_36_factsheets,  # 40
    make_slide_37_threat_model,  # 41
    make_slide_38_annotated3,  # 42
    make_slide_39_ibm_strengths_gaps,  # 43
    # Block 4 — Hands-on (40-45)
    make_slide_40_divider_handson,  # 44
    make_slide_41_setup,  # 45
    make_slide_42_notebook_map,  # 46
    make_slide_43_live_nb3,  # 47
    make_slide_44_live_nb6,  # 48
    make_slide_45_critique,  # 49
    # Block 5 — Close (46-49)
    make_slide_46_takeaway,  # 50
    make_slide_47_monday_actions,  # 51
    make_slide_48_pointers,  # 52
    make_slide_49_closing,  # 53
]


def build_deck() -> Path:
    """Build the complete 53-slide presentation and write to disk."""
    prs = Presentation()

    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for builder in SLIDE_BUILDERS:
        builder(prs)

    out_dir = Path(__file__).parent
    out_path = out_dir / "data-architecture-ai-era.pptx"
    prs.save(str(out_path))

    slide_count = len(prs.slides)
    file_size = out_path.stat().st_size
    print(f"Deck saved: {out_path}")
    print(f"  Slides : {slide_count}")
    print(f"  Size   : {file_size:,} bytes ({file_size / 1024:.0f} KB)")

    if slide_count != 55:
        print(f"  WARNING: expected 55 slides, got {slide_count}")
    if file_size < 100_000:
        print(f"  WARNING: file is under 100 KB ({file_size:,} bytes)")

    return out_path


if __name__ == "__main__":
    build_deck()
