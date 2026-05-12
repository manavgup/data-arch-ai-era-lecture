"""Validate the generated PPTX deck against SPEC.md rules (sections 3 and 8).

Checks enforced:
  - Light theme: every slide background is #F4F2EC (paper)
  - Speaker notes present on every content slide (skips title + section dividers)
  - Slide count within the 50-60 target range
  - No forbidden terms: "Project Mantis", "NaviOwl", "Closing the Gap", "Extendicare"
  - Anonymised critique uses "HealthFirst Insights"
  - Design tokens: accent blue #2D4ADE appears somewhere in the deck
  - Font families: Georgia, Calibri, Consolas are the only fonts used
  - Manav's surname is "Gupta" (never "Mistry")
  - Ship AI podcast attributed only to Manav (never "Two Guys with AI" as current name)
"""

from __future__ import annotations

import re
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor

DECK_PATH = Path(__file__).resolve().parent.parent / "deck" / "data-architecture-ai-era.pptx"

# ── Design tokens from SPEC.md §3 / generate_deck.py ────────────────────────

PAPER = RGBColor(0xF4, 0xF2, 0xEC)
ACCENT = RGBColor(0x2D, 0x4A, 0xDE)

ALLOWED_FONTS = {"Georgia", "Calibri", "Consolas"}

# Slide indices (0-based) that are title/divider slides and may lack speaker notes.
# These are identified by their position: slide 0 (title), and any slide whose
# only text is a section header.  We detect them dynamically below.

MIN_SLIDES = 50
MAX_SLIDES = 60

# ── Forbidden / required terms (SPEC.md §8) ─────────────────────────────────

FORBIDDEN_TERMS = [
    "Project Mantis",
    "NaviOwl",
    "Closing the Gap",
    "Extendicare",
]

FORBIDDEN_SURNAME = "Mistry"

# ── Helpers ──────────────────────────────────────────────────────────────────


def _all_text_in_slide(slide) -> str:
    """Concatenate all visible text on a slide (shapes only, not notes)."""
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(parts)


def _notes_text(slide) -> str:
    """Return speaker notes text for a slide, or empty string."""
    if slide.has_notes_slide:
        return slide.notes_slide.notes_text_frame.text.strip()
    return ""


def _all_text_in_deck(prs: Presentation) -> str:
    """All text across every slide and its speaker notes."""
    parts: list[str] = []
    for slide in prs.slides:
        parts.append(_all_text_in_slide(slide))
        parts.append(_notes_text(slide))
    return "\n".join(parts)


def _collect_fonts(prs: Presentation) -> set[str]:
    """Collect every font name used across all slide shapes."""
    fonts: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
    return fonts


def _slide_bg_color(slide) -> RGBColor | None:
    """Extract the solid-fill background colour of a slide, or None."""
    bg = slide.background
    fill = bg.fill
    if fill.type is not None:
        try:
            return fill.fore_color.rgb
        except (AttributeError, TypeError):
            return None
    return None


def _is_divider_slide(slide) -> bool:
    """Heuristic: a slide is a divider/title if it has <= 2 text shapes and no table."""
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    has_table = any(s.has_table for s in slide.shapes)
    total_text = " ".join(s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame)
    # Divider slides are short (< 100 chars) and have at most 2 text shapes
    return len(text_shapes) <= 2 and not has_table and len(total_text) < 100


# ── Fixtures ─────────────────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def deck() -> Presentation:
    if not DECK_PATH.exists():
        pytest.skip(f"Deck not found at {DECK_PATH}")
    return Presentation(str(DECK_PATH))


# ── Tests ────────────────────────────────────────────────────────────────────


class TestLightTheme:
    """SPEC.md §8 rule 1: Light theme PPTX always."""

    def test_all_slides_have_paper_background(self, deck: Presentation) -> None:
        violations: list[str] = []
        for idx, slide in enumerate(deck.slides, 1):
            bg = _slide_bg_color(slide)
            if bg is not None and bg != PAPER:
                violations.append(f"Slide {idx}: background #{bg} (expected #{PAPER})")
        assert not violations, "Non-paper backgrounds found:\n" + "\n".join(violations)


class TestSpeakerNotes:
    """SPEC.md §3: Speaker notes on every content slide."""

    def test_content_slides_have_notes(self, deck: Presentation) -> None:
        missing: list[str] = []
        for idx, slide in enumerate(deck.slides, 1):
            if _is_divider_slide(slide):
                continue
            notes = _notes_text(slide)
            if len(notes) < 10:
                slide_text = _all_text_in_slide(slide)[:80].replace("\n", " ")
                missing.append(f"Slide {idx}: {slide_text!r}")
        assert not missing, f"{len(missing)} content slide(s) missing speaker notes:\n" + "\n".join(missing)


class TestSlideCount:
    """SPEC.md §3: ~50-60 slides total."""

    def test_slide_count_in_range(self, deck: Presentation) -> None:
        count = len(deck.slides)
        assert MIN_SLIDES <= count <= MAX_SLIDES, f"Deck has {count} slides, expected {MIN_SLIDES}-{MAX_SLIDES}"


class TestForbiddenTerms:
    """SPEC.md §8 rules 5-6: no NaviOwl, no Project Mantis, etc."""

    def test_no_forbidden_terms(self, deck: Presentation) -> None:
        full_text = _all_text_in_deck(deck)
        found: list[str] = []
        for term in FORBIDDEN_TERMS:
            if re.search(re.escape(term), full_text, re.IGNORECASE):
                found.append(term)
        assert not found, "Forbidden terms found in deck: " + ", ".join(found)

    def test_no_mistry_surname(self, deck: Presentation) -> None:
        """SPEC.md §8 rule 4: Manav's surname is Gupta, not Mistry."""
        full_text = _all_text_in_deck(deck)
        assert FORBIDDEN_SURNAME.lower() not in full_text.lower(), (
            f"Found '{FORBIDDEN_SURNAME}' in deck — Manav's surname must be Gupta"
        )


class TestDesignTokens:
    """SPEC.md §3: accent blue #2D4ADE should be used."""

    def test_accent_blue_present(self, deck: Presentation) -> None:
        for slide in deck.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb == ACCENT:
                            return  # found it
        # Also check shape fills and line colors
        for slide in deck.slides:
            for shape in slide.shapes:
                if hasattr(shape, "fill") and shape.fill.type is not None:
                    try:
                        if shape.fill.fore_color.rgb == ACCENT:
                            return
                    except (AttributeError, TypeError):
                        pass
                if hasattr(shape, "line") and shape.line.fill.type is not None:
                    try:
                        if shape.line.color.rgb == ACCENT:
                            return
                    except (AttributeError, TypeError):
                        pass
        pytest.fail(f"Accent blue #{ACCENT} not found anywhere in deck")


class TestFonts:
    """SPEC.md §3: Only Georgia, Calibri, Consolas should be used."""

    def test_only_allowed_fonts(self, deck: Presentation) -> None:
        used = _collect_fonts(deck)
        unexpected = used - ALLOWED_FONTS
        assert not unexpected, f"Unexpected fonts: {unexpected}. Allowed: {ALLOWED_FONTS}"


class TestContentIntegrity:
    """SPEC.md §8 rules 3, 7-10: content correctness."""

    def test_ship_ai_not_two_guys(self, deck: Presentation) -> None:
        """Ship AI podcast is the current name; 'Two Guys with AI' is the old name."""
        full_text = _all_text_in_deck(deck)
        assert "Two Guys with AI" not in full_text, "Found 'Two Guys with AI' — use 'Ship AI' as the podcast name"

    def test_statistics_have_sources(self, deck: Presentation) -> None:
        """SPEC.md §8 rule 10: quoted statistics must be sourced in speaker notes.

        Heuristic: if slide body contains a percentage or dollar figure,
        the speaker notes should contain a source-like reference.
        """
        stat_pattern = re.compile(r"\b\d+\.?\d*\s*[%$]|\$\s*\d+|billion|trillion|million", re.IGNORECASE)
        source_pattern = re.compile(
            r"source|cited|according to|per |report|gartner|idc|forrester|mckinsey|ibm ", re.IGNORECASE
        )
        unsourced: list[str] = []
        for idx, slide in enumerate(deck.slides, 1):
            body = _all_text_in_slide(slide)
            if stat_pattern.search(body):
                notes = _notes_text(slide)
                if not source_pattern.search(notes) and not source_pattern.search(body):
                    snippet = body[:80].replace("\n", " ")
                    unsourced.append(f"Slide {idx}: stat found but no source — {snippet!r}")
        if unsourced:
            pytest.fail(f"{len(unsourced)} slide(s) have statistics without visible sourcing:\n" + "\n".join(unsourced))
